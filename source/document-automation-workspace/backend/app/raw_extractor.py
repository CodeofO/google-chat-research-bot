import base64
import html
import json
import mimetypes
import os
import platform
import shutil
import subprocess
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Any

try:
    import pymupdf as fitz
except ImportError:  # pragma: no cover - compatibility for older PyMuPDF installs
    import fitz
from fastapi import UploadFile

from app.config import get_settings


RAW_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".pdf"}
OFFICE_PDF_FILTERS = {
    ".docx": "writer_pdf_Export",
    ".xlsx": "calc_pdf_Export",
    ".pptx": "impress_pdf_Export",
}
PDF_EXPORT_OPTIONS = {
    "MaxImageResolution": {"type": "long", "value": "300"},
    "Quality": {"type": "long", "value": "95"},
    "ReduceImageResolution": {"type": "boolean", "value": "false"},
    "EmbedStandardFonts": {"type": "boolean", "value": "true"},
    "EmbedFonts": {"type": "boolean", "value": "true"},
    "SubsetFonts": {"type": "boolean", "value": "false"},
}
OFFICE_MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
IMAGE_EXTENSIONS = {".gif", ".jpeg", ".jpg", ".png", ".svg", ".webp"}


@dataclass(frozen=True)
class RawExtractionOptions:
    include_images: bool = True
    include_formulas: bool = False


class RawExtractionError(ValueError):
    def __init__(self, message: str, status_code: int = 400):
        self.status_code = status_code
        super().__init__(message)


def validate_raw_upload(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in RAW_EXTENSIONS:
        raise RawExtractionError("Only DOCX, XLSX, PPTX, and PDF files are supported")
    return suffix


def save_raw_upload(upload: UploadFile, raw_id: str) -> tuple[str, str, Path, int]:
    suffix = validate_raw_upload(upload.filename or "")
    settings = get_settings()
    raw_dir = settings.resolved_raw_storage_dir / raw_id
    raw_dir.mkdir(parents=True, exist_ok=True)
    original_path = raw_dir / f"original{suffix}"

    size = 0
    with original_path.open("wb") as destination:
        while True:
            chunk = upload.file.read(1024 * 1024)
            if not chunk:
                break
            size += len(chunk)
            if size > settings.upload_max_file_bytes:
                raise RawExtractionError("Uploaded file is too large", status_code=413)
            destination.write(chunk)

    _validate_raw_file_content(original_path, suffix)
    return upload.filename or original_path.name, suffix[1:], original_path, size


def _validate_raw_file_content(path: Path, suffix: str) -> None:
    if suffix == ".pdf":
        if not path.read_bytes()[:16].startswith(b"%PDF-"):
            raise RawExtractionError("Uploaded file does not match the PDF format", status_code=415)
        return
    if suffix in {".docx", ".xlsx", ".pptx"}:
        if not zipfile.is_zipfile(path):
            raise RawExtractionError("Uploaded file does not match the Office document format", status_code=415)
        try:
            with zipfile.ZipFile(path) as archive:
                names = set(archive.namelist())
        except zipfile.BadZipFile as exc:
            raise RawExtractionError("Uploaded file does not match the Office document format", status_code=415) from exc
        if "[Content_Types].xml" not in names:
            raise RawExtractionError("Uploaded Office document is missing required metadata", status_code=415)
        expected_dir = {".docx": "word/", ".xlsx": "xl/", ".pptx": "ppt/"}[suffix]
        if not any(name.startswith(expected_dir) for name in names):
            raise RawExtractionError(f"Uploaded file does not match the {suffix[1:].upper()} format", status_code=415)


def create_raw_outputs(
    source_path: Path,
    source_format: str,
    options: RawExtractionOptions | None = None,
) -> tuple[Path, Path, list[str]]:
    options = options or RawExtractionOptions()
    suffix = f".{source_format.lower()}"
    output_dir = source_path.parent
    pdf_path = output_dir / "preview.pdf"
    html_path = output_dir / "content.html"
    warnings: list[str] = []

    create_pdf_preview(source_path, suffix, pdf_path)
    html_path.write_text(build_html_document(source_path, suffix, options), encoding="utf-8")
    if html_path.stat().st_size == 0:
        warnings.append("empty_html_output")
    return pdf_path, html_path, warnings


def create_pdf_preview(source_path: Path, suffix: str, pdf_path: Path) -> None:
    if suffix == ".pdf":
        shutil.copy2(source_path, pdf_path)
        return
    convert_office_to_pdf(source_path, suffix, pdf_path)


def convert_office_to_pdf(source_path: Path, suffix: str, pdf_path: Path) -> None:
    export_filter = OFFICE_PDF_FILTERS.get(suffix)
    if not export_filter:
        raise RawExtractionError(f"PDF conversion is not supported for {suffix}")

    soffice = find_libreoffice()
    with tempfile.TemporaryDirectory(prefix="raw2pdf_") as tmp:
        tmp_path = Path(tmp)
        cmd = [
            str(soffice),
            "--headless",
            "--convert-to",
            f"pdf:{export_filter}:{json.dumps(PDF_EXPORT_OPTIONS)}",
            "--outdir",
            str(tmp_path),
            str(source_path),
        ]
        completed = subprocess.run(cmd, capture_output=True, text=True, timeout=600, env=libreoffice_environment())
        if completed.returncode != 0:
            raise RawExtractionError(
                f"LibreOffice failed to convert the document to PDF: {completed.stderr or completed.stdout}"
            )
        produced = tmp_path / f"{source_path.stem}.pdf"
        if not produced.is_file():
            raise RawExtractionError("LibreOffice did not produce a PDF preview")
        shutil.move(str(produced), pdf_path)


def find_libreoffice() -> Path:
    settings = get_settings()
    candidates: list[Path] = []
    if settings.libreoffice_path:
        candidates.append(Path(settings.libreoffice_path))
    if platform.system() == "Darwin":
        candidates.append(Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"))
    candidates.extend(
        [
            Path("/usr/bin/soffice"),
            Path("/usr/local/bin/soffice"),
            Path("/opt/homebrew/bin/soffice"),
            Path("/opt/libreoffice25.2/program/soffice"),
        ]
    )
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    raise RawExtractionError(
        "LibreOffice was not found. Install LibreOffice or set LIBREOFFICE_PATH in the backend .env."
    )


def libreoffice_environment() -> dict[str, str]:
    env = os.environ.copy()
    if platform.system() == "Darwin":
        env["PYTHONPATH"] = (
            "/Applications/LibreOffice.app/Contents/Frameworks/"
            "LibreOfficePython.framework/Versions/3.10/lib/python3.10"
        )
        env["PYTHONHOME"] = "/Applications/LibreOffice.app/Contents/Frameworks/LibreOfficePython.framework/Versions/3.10"
    elif platform.system() == "Linux":
        env["PYTHONPATH"] = "/opt/libreoffice25.2/program"
        env["PYTHONHOME"] = "/opt/libreoffice25.2/program/python"
    return env


def build_html_document(
    source_path: Path,
    suffix: str,
    options: RawExtractionOptions | None = None,
) -> str:
    options = options or RawExtractionOptions()
    if suffix == ".docx":
        body = _docx_to_html(source_path, options)
    elif suffix == ".xlsx":
        body = _xlsx_to_html(source_path, options)
    elif suffix == ".pptx":
        body = _pptx_to_html(source_path, options)
    elif suffix == ".pdf":
        body = _pdf_to_html(source_path, options)
    else:
        raise RawExtractionError(f"HTML extraction is not supported for {suffix}")

    title = html.escape(source_path.name)
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{title}</title>
  <style>
    body {{ color: #1d2529; font-family: Inter, Arial, sans-serif; line-height: 1.55; margin: 24px; }}
    h1, h2, h3 {{ line-height: 1.2; }}
    section {{ border-bottom: 1px solid #dde4df; margin-bottom: 24px; padding-bottom: 18px; }}
    table {{ border-collapse: collapse; margin: 12px 0; width: 100%; }}
    td, th {{ border: 1px solid #cbd5cf; padding: 6px 8px; vertical-align: top; }}
    th {{ background: #eef3ef; }}
    figure {{ margin: 16px 0; }}
    figcaption {{ color: #53615a; font-size: 12px; margin-top: 4px; }}
    img {{ border: 1px solid #d8e0da; border-radius: 6px; height: auto; max-width: 100%; }}
    .formula {{ background: #eef3ef; border-radius: 4px; color: #0f625e; display: inline-block; font-family: monospace; padding: 2px 4px; }}
    .formula-value {{ color: #53615a; display: block; font-size: 12px; margin-top: 3px; }}
    pre {{ white-space: pre-wrap; }}
  </style>
</head>
<body>
  <h1>{title}</h1>
  {body}
</body>
</html>"""


def _docx_to_html(source_path: Path, options: RawExtractionOptions) -> str:
    import bleach
    import mammoth

    with source_path.open("rb") as source:
        result = mammoth.convert_to_html(
            source,
            convert_image=mammoth.images.data_uri if options.include_images else _ignore_mammoth_image,
        )
    allowed_tags = {
        "a",
        "b",
        "br",
        "em",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "i",
        "li",
        "ol",
        "p",
        "strong",
        "table",
        "tbody",
        "td",
        "th",
        "thead",
        "tr",
        "u",
        "ul",
    }
    if options.include_images:
        allowed_tags.add("img")
    allowed_attrs = {"a": ["href"], "td": ["colspan", "rowspan"], "th": ["colspan", "rowspan"]}
    if options.include_images:
        allowed_attrs["img"] = ["alt", "src"]
    protocols = bleach.sanitizer.ALLOWED_PROTOCOLS | {"data"}
    content = bleach.clean(result.value, tags=allowed_tags, attributes=allowed_attrs, protocols=protocols, strip=True)
    if options.include_images:
        fallback_images = _embedded_package_images_section(source_path, "Document images", ["word/media/"])
        if fallback_images:
            content += fallback_images
    if options.include_formulas:
        formulas = _extract_ooxml_formulas(source_path, ["word/document.xml"])
        if formulas:
            content += _formula_section("Formulas", formulas)
    return content


def _xlsx_to_html(source_path: Path, options: RawExtractionOptions) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(source_path, read_only=True, data_only=True)
    formula_workbook = load_workbook(source_path, read_only=True, data_only=False) if options.include_formulas else None
    sections: list[str] = []
    try:
        for sheet in workbook.worksheets:
            if formula_workbook:
                formula_sheet = formula_workbook[sheet.title]
                rows = _worksheet_rows_with_formulas(
                    list(sheet.iter_rows(values_only=True)),
                    list(formula_sheet.iter_rows(values_only=True)),
                )
            else:
                rows = list(sheet.iter_rows(values_only=True))
            rows = [row for row in rows if any(_cell_has_value(cell) for cell in row)]
            table = _rows_to_table(rows)
            sections.append(f"<section><h2>{html.escape(sheet.title)}</h2>{table}</section>")
        if options.include_images:
            images = _embedded_package_images_section(source_path, "Workbook images", ["xl/media/"])
            if images:
                sections.append(images)
    finally:
        workbook.close()
        if formula_workbook:
            formula_workbook.close()
    return "\n".join(sections) or "<p>No readable worksheet data found.</p>"


def _pptx_to_html(source_path: Path, options: RawExtractionOptions) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(source_path)
    formulas_by_slide = _extract_pptx_formulas_by_slide(source_path) if options.include_formulas else {}
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        fragments = [f"<h2>Slide {index}</h2>"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _shape_text_to_html(shape)
                if text:
                    fragments.append(text)
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                fragments.append(_rows_to_table(rows))
            if options.include_images and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image_html = _pptx_picture_to_html(shape, f"Slide {index} image")
                if image_html:
                    fragments.append(image_html)
        formulas = formulas_by_slide.get(index, [])
        if formulas:
            fragments.append(_formula_list("Slide formulas", formulas, "h3"))
        sections.append(f"<section>{''.join(fragments)}</section>")
    if options.include_images:
        fallback_images = _embedded_package_images_section(source_path, "Presentation images", ["ppt/media/"])
        if fallback_images:
            sections.append(fallback_images)
    return "\n".join(sections) or "<p>No readable slide content found.</p>"


def _pdf_to_html(source_path: Path, options: RawExtractionOptions) -> str:
    sections: list[str] = []
    with fitz.open(source_path) as document:
        for page_index, page in enumerate(document, start=1):
            blocks = page.get_text("blocks")
            fragments = [f"<h2>Page {page_index}</h2>"]
            for block in sorted(blocks, key=lambda item: (item[1], item[0])):
                text = str(block[4]).strip()
                if text:
                    fragments.append(f"<p>{html.escape(text).replace(chr(10), '<br />')}</p>")
            if options.include_images:
                for image_index, image_info in enumerate(page.get_images(full=True), start=1):
                    extracted = document.extract_image(image_info[0])
                    image = _image_data_uri_html(
                        extracted["image"],
                        _image_content_type(extracted.get("ext")),
                        f"Page {page_index} image {image_index}",
                    )
                    if image:
                        fragments.append(image)
            sections.append(f"<section>{''.join(fragments)}</section>")
    return "\n".join(sections) or "<p>No readable PDF text found.</p>"


def _shape_text_to_html(shape: Any) -> str:
    paragraphs: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            paragraphs.append(f"<p>{html.escape(text)}</p>")
    return "".join(paragraphs)


def _rows_to_table(rows: list[Any]) -> str:
    if not rows:
        return "<p>No table data found.</p>"
    html_rows: list[str] = []
    for row_index, row in enumerate(rows):
        tag = "th" if row_index == 0 else "td"
        cells = "".join(f"<{tag}>{_cell_to_html(cell)}</{tag}>" for cell in row)
        html_rows.append(f"<tr>{cells}</tr>")
    return f"<table><tbody>{''.join(html_rows)}</tbody></table>"


@dataclass(frozen=True)
class FormulaCell:
    formula: str
    cached_value: Any


def _worksheet_rows_with_formulas(data_rows: list[Any], formula_rows: list[Any]) -> list[list[Any]]:
    merged_rows: list[list[Any]] = []
    row_count = max(len(data_rows), len(formula_rows))
    for row_index in range(row_count):
        data_row = data_rows[row_index] if row_index < len(data_rows) else ()
        formula_row = formula_rows[row_index] if row_index < len(formula_rows) else ()
        column_count = max(len(data_row), len(formula_row))
        merged_row: list[Any] = []
        for column_index in range(column_count):
            data_value = data_row[column_index] if column_index < len(data_row) else None
            formula_value = formula_row[column_index] if column_index < len(formula_row) else None
            if isinstance(formula_value, str) and formula_value.startswith("="):
                merged_row.append(FormulaCell(formula=formula_value, cached_value=data_value))
            else:
                merged_row.append(data_value)
        merged_rows.append(merged_row)
    return merged_rows


def _cell_has_value(value: Any) -> bool:
    if isinstance(value, FormulaCell):
        return bool(value.formula)
    return value is not None


def _cell_to_html(value: Any) -> str:
    if isinstance(value, FormulaCell):
        formula = html.escape(value.formula)
        if value.cached_value not in (None, ""):
            cached = html.escape(_cell_to_text(value.cached_value))
            return f'<span class="formula">{formula}</span><span class="formula-value">cached value: {cached}</span>'
        return f'<span class="formula">{formula}</span>'
    return html.escape(_cell_to_text(value))


def _cell_to_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value)


def _ignore_mammoth_image(_: Any) -> list[Any]:
    return []


def _pptx_picture_to_html(shape: Any, fallback_alt: str) -> str:
    try:
        image = shape.image
        alt = getattr(shape, "name", None) or fallback_alt
        return _image_data_uri_html(image.blob, image.content_type, alt)
    except Exception:
        return ""


def _embedded_package_images_section(source_path: Path, title: str, prefixes: list[str]) -> str:
    figures: list[str] = []
    try:
        with zipfile.ZipFile(source_path) as archive:
            for member in archive.namelist():
                if not any(member.startswith(prefix) for prefix in prefixes):
                    continue
                suffix = Path(member).suffix.lower()
                if suffix not in IMAGE_EXTENSIONS:
                    continue
                content_type = mimetypes.guess_type(member)[0] or "application/octet-stream"
                image = _image_data_uri_html(archive.read(member), content_type, Path(member).name)
                if image:
                    figures.append(image)
    except zipfile.BadZipFile:
        return ""
    if not figures:
        return ""
    return f"<section><h2>{html.escape(title)}</h2>{''.join(figures)}</section>"


def _image_data_uri_html(blob: bytes, content_type: str | None, alt: str) -> str:
    if not blob:
        return ""
    content_type = content_type or "application/octet-stream"
    if content_type not in {"image/gif", "image/jpeg", "image/png", "image/svg+xml", "image/webp"}:
        return ""
    encoded = base64.b64encode(blob).decode("ascii")
    escaped_alt = html.escape(alt, quote=True)
    return (
        f'<figure><img src="data:{content_type};base64,{encoded}" alt="{escaped_alt}" />'
        f"<figcaption>{escaped_alt}</figcaption></figure>"
    )


def _image_content_type(extension: str | None) -> str | None:
    if not extension:
        return None
    extension = extension.lower().lstrip(".")
    if extension == "jpg":
        extension = "jpeg"
    if extension == "svg":
        return "image/svg+xml"
    return f"image/{extension}"


def _formula_section(title: str, formulas: list[str]) -> str:
    return f"<section>{_formula_list(title, formulas, 'h2')}</section>"


def _formula_list(title: str, formulas: list[str], heading_tag: str) -> str:
    if not formulas:
        return ""
    items = "".join(f'<li><span class="formula">{html.escape(formula)}</span></li>' for formula in formulas)
    return f"<{heading_tag}>{html.escape(title)}</{heading_tag}><ol>{items}</ol>"


def _extract_ooxml_formulas(source_path: Path, members: list[str]) -> list[str]:
    formulas: list[str] = []
    try:
        with zipfile.ZipFile(source_path) as archive:
            names = set(archive.namelist())
            for member in members:
                if member not in names:
                    continue
                formulas.extend(_formulas_from_xml(archive.read(member)))
    except (zipfile.BadZipFile, ET.ParseError):
        return []
    return _dedupe(formulas)


def _extract_pptx_formulas_by_slide(source_path: Path) -> dict[int, list[str]]:
    formulas_by_slide: dict[int, list[str]] = {}
    try:
        with zipfile.ZipFile(source_path) as archive:
            slide_names = sorted(
                (name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")),
                key=_slide_sort_key,
            )
            for index, member in enumerate(slide_names, start=1):
                formulas = _dedupe(_formulas_from_xml(archive.read(member)))
                if formulas:
                    formulas_by_slide[index] = formulas
    except (zipfile.BadZipFile, ET.ParseError):
        return {}
    return formulas_by_slide


def _formulas_from_xml(xml_bytes: bytes) -> list[str]:
    root = ET.fromstring(xml_bytes)
    formulas: list[str] = []
    for element in root.iter(f"{{{OFFICE_MATH_NS}}}oMath"):
        text = "".join(node.text or "" for node in element.iter(f"{{{OFFICE_MATH_NS}}}t")).strip()
        if text:
            formulas.append(text)
    return formulas


def _dedupe(values: list[str]) -> list[str]:
    seen: set[str] = set()
    deduped: list[str] = []
    for value in values:
        normalized = " ".join(value.split())
        if normalized and normalized not in seen:
            seen.add(normalized)
            deduped.append(normalized)
    return deduped


def _slide_sort_key(name: str) -> tuple[int, str]:
    stem = Path(name).stem
    index = stem.replace("slide", "")
    return (int(index) if index.isdigit() else 0, name)
