import base64
import asyncio
import io
import json
import os
import re
import threading
import time
from datetime import datetime, timedelta
from pathlib import Path
from types import SimpleNamespace
from typing import Any
from uuid import uuid4
import zipfile

import pytest

try:
    import pymupdf as fitz
except ImportError:  # pragma: no cover - compatibility for older PyMuPDF installs
    import fitz
from PIL import Image

from app.config import get_settings
from app.models import (
    AuditEvent,
    Batch,
    ClassificationBatch,
    Document,
    ExportJob,
    RawExtraction,
    RequiredFieldCheckBatch,
    WorkflowRun,
    WorkflowRunItem,
)
from tests.conftest import get_client


ONE_BY_ONE_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)
SCHEMA_COUNTER = 0
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def assert_xlsx_response(response) -> None:
    assert response.status_code == 200, response.text
    assert response.headers["content-type"] == XLSX_MIME
    assert response.content.startswith(b"PK")


def wait_for_export_job(client, job_id: str) -> dict[str, Any]:
    for _ in range(20):
        response = client.get(f"/api/export-jobs/{job_id}")
        assert response.status_code == 200, response.text
        payload = response.json()
        if payload["status"] in {"completed", "failed"}:
            return payload
        time.sleep(0.01)
    raise AssertionError(f"Export job did not finish: {job_id}")


def _read_env_example(path: Path) -> dict[str, str]:
    values: dict[str, str] = {}
    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        values[key.strip()] = value.strip()
    return values


def test_health() -> None:
    with get_client() as client:
        response = client.get("/api/health")
    assert response.status_code == 200
    assert response.json() == {"status": "ok"}


def test_system_status_mock_mode(monkeypatch) -> None:
    try:
        monkeypatch.setenv("VLM_PROVIDER", "mock")
        monkeypatch.setenv("UPLOAD_MAX_BATCH_FILES", "1234")
        monkeypatch.setenv("UPLOAD_CHUNK_FILES", "10")
        monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "16")
        monkeypatch.setenv("VLM_MAX_CONCURRENT_REQUESTS", "8")
        get_settings.cache_clear()
        with get_client() as client:
            response = client.get("/api/system/status")
        assert response.status_code == 200
        payload = response.json()
        assert payload["vlm_provider"] == "mock"
        assert payload["is_mock"] is True
        assert payload["upload_max_batch_files"] == 1234
        assert payload["upload_chunk_files"] == 10
        assert payload["workflow_max_workers"] == 16
        assert payload["vlm_max_concurrent_requests"] == 8
        assert "batch_max_workers" not in payload
        assert payload["document_page_max_long_edge"] == 3000
        assert payload["document_page_jpeg_quality"] == 88
        assert "vlm_api_key" not in payload
    finally:
        get_settings.cache_clear()


def test_database_pool_defaults_to_64_connections() -> None:
    from app.database import engine

    pool_size = getattr(engine.pool, "size", None)
    assert callable(pool_size)
    assert pool_size() == 64
    assert getattr(engine.pool, "_max_overflow", None) == 0


def test_run_workflow_blocking_queues_without_threadpool_deadlock(monkeypatch) -> None:
    from app.concurrency import run_workflow_blocking

    active = 0
    max_active = 0
    lock = threading.Lock()

    def blocking_work(index: int) -> int:
        nonlocal active, max_active
        with lock:
            active += 1
            max_active = max(max_active, active)
        time.sleep(0.01)
        with lock:
            active -= 1
        return index

    async def run_many() -> list[int]:
        tasks = [run_workflow_blocking(blocking_work, index) for index in range(80)]
        return await asyncio.wait_for(asyncio.gather(*tasks), timeout=3)

    try:
        monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
        get_settings.cache_clear()
        assert asyncio.run(run_many()) == list(range(80))
        assert max_active == 2
    finally:
        get_settings.cache_clear()


def test_gather_workflow_limited_caps_async_tasks(monkeypatch) -> None:
    from app.concurrency import gather_workflow_limited

    active = 0
    max_active = 0

    async def async_work(index: int) -> int:
        nonlocal active, max_active
        active += 1
        max_active = max(max_active, active)
        await asyncio.sleep(0.01)
        active -= 1
        return index

    try:
        monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
        get_settings.cache_clear()
        assert asyncio.run(gather_workflow_limited(list(range(6)), async_work)) == list(range(6))
        assert max_active == 2
    finally:
        get_settings.cache_clear()


def test_async_vlm_limit_queues_without_threadpool_deadlock(monkeypatch) -> None:
    from app import vlm as vlm_module

    active = 0
    max_active = 0
    lock = asyncio.Lock()

    async def fake_invoke(*args, **kwargs):
        nonlocal active, max_active
        async with lock:
            active += 1
            max_active = max(max_active, active)
        await asyncio.sleep(0.01)
        async with lock:
            active -= 1
        return {"ok": True}

    async def run_many() -> list[dict[str, bool]]:
        tasks = [
            vlm_module._invoke_vlm_with_limit_async("system", "prompt", [], {}, "google_genai")
            for _ in range(80)
        ]
        return await asyncio.wait_for(asyncio.gather(*tasks), timeout=3)

    try:
        monkeypatch.setenv("VLM_MAX_CONCURRENT_REQUESTS", "2")
        monkeypatch.setenv("VLM_TIMEOUT_SECONDS", "2")
        monkeypatch.setenv("VLM_MAX_RETRIES", "0")
        get_settings.cache_clear()
        monkeypatch.setattr(vlm_module, "_invoke_structured_llm_async", fake_invoke)

        results = asyncio.run(run_many())
        assert results == [{"ok": True}] * 80
        assert max_active == 2
        assert vlm_module.vlm_runtime_counters()["vlm_active_count"] == 0
        assert vlm_module.vlm_runtime_counters()["vlm_waiting_count"] == 0
    finally:
        get_settings.cache_clear()


def test_async_vlm_request_times_out(monkeypatch) -> None:
    from app import vlm as vlm_module

    async def slow_invoke(*args, **kwargs):
        await asyncio.sleep(2)
        return {"ok": True}

    async def run_once() -> None:
        with pytest.raises(vlm_module.VlmRuntimeError) as exc_info:
            await vlm_module._invoke_vlm_with_limit_async("system", "prompt", [], {}, "google_genai")
        assert "timed out" in exc_info.value.message

    try:
        monkeypatch.setenv("VLM_MAX_CONCURRENT_REQUESTS", "2")
        monkeypatch.setenv("VLM_TIMEOUT_SECONDS", "1")
        monkeypatch.setenv("VLM_MAX_RETRIES", "0")
        get_settings.cache_clear()
        monkeypatch.setattr(vlm_module, "_invoke_structured_llm_async", slow_invoke)

        asyncio.run(run_once())
        assert vlm_module.vlm_runtime_counters()["vlm_active_count"] == 0
    finally:
        get_settings.cache_clear()


def test_root_env_upsert_creates_vlm_settings(monkeypatch, tmp_path) -> None:
    from app import config as config_module

    env_path = tmp_path / ".env"
    monkeypatch.setattr(config_module, "ROOT_ENV_PATH", env_path)
    config_module.upsert_root_env(
        {
            "VLM_API_KEY": "test-secret",
            "VLM_MODEL_NAME": "test-model",
            "LIBREOFFICE_PATH": "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        },
        include_defaults=True,
    )

    contents = env_path.read_text(encoding="utf-8")
    assert 'APP_ENV="local"' in contents
    assert 'VLM_API_KEY="test-secret"' in contents
    assert 'VLM_MODEL_NAME="test-model"' in contents
    assert 'LIBREOFFICE_PATH="/Applications/LibreOffice.app/Contents/MacOS/soffice"' in contents


def test_vlm_settings_include_libreoffice_path(monkeypatch, tmp_path) -> None:
    from app import config as config_module

    env_path = tmp_path / ".env"
    monkeypatch.setattr(config_module, "ROOT_ENV_PATH", env_path)

    with get_client() as client:
        response = client.put(
            "/api/settings/vlm",
            json={
                "api_key": "test-secret",
                "model_name": "test-model",
                "base_url": "http://127.0.0.1:11434/v1",
                "libreoffice_path": "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "provider": "openai",
                "workflow_max_workers": 11,
                "vlm_max_concurrent_requests": 7,
                "vlm_timeout_seconds": 600,
                "kie_field_group_size": 3,
                "inference_params": {
                    "reasoning_effort": "off",
                    "thinking": "off",
                    "temperature": "0",
                    "verbosity": "",
                    "max_completion_tokens": "2048",
                    "top_p": "0.9",
                    "service_tier": "",
                },
            },
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["base_url"] == "http://127.0.0.1:11434/v1"
        assert payload["libreoffice_path"] == "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        assert payload["inference_params"]["reasoning_effort"] == "off"
        assert payload["inference_params"]["thinking"] == "off"
        assert payload["inference_params"]["max_completion_tokens"] == "2048"
        assert "vlm_max_concurrent_requests" in payload
        assert payload["workflow_max_workers"] == 11
        assert payload["vlm_timeout_seconds"] == 600
        assert "batch_max_workers" not in payload
        assert "kie_field_group_size" in payload
        assert payload["runtime_settings_writable"] is True

    contents = env_path.read_text(encoding="utf-8")
    assert 'VLM_BASE_URL="http://127.0.0.1:11434/v1"' in contents
    assert 'LIBREOFFICE_PATH="/Applications/LibreOffice.app/Contents/MacOS/soffice"' in contents
    assert "BATCH_MAX_WORKERS" not in contents
    assert 'WORKFLOW_MAX_WORKERS="11"' in contents
    assert 'VLM_MAX_CONCURRENT_REQUESTS="7"' in contents
    assert 'VLM_TIMEOUT_SECONDS="600"' in contents
    assert 'KIE_FIELD_GROUP_SIZE="3"' in contents
    assert "VLM_INFERENCE_PARAMS=" in contents
    assert "VLM_TEMPERATURE" not in contents
    assert "VLM_REASONING_EFFORT" not in contents
    assert "VLM_VERBOSITY" not in contents


def test_vlm_settings_are_readonly_in_production(monkeypatch, tmp_path) -> None:
    from app import config as config_module

    env_path = tmp_path / ".env"
    monkeypatch.setattr(config_module, "ROOT_ENV_PATH", env_path)
    monkeypatch.setenv("APP_ENV", "production")
    monkeypatch.delenv("ALLOW_RUNTIME_SETTINGS", raising=False)
    get_settings.cache_clear()

    try:
        with get_client() as client:
            settings_response = client.get("/api/settings/vlm")
            assert settings_response.status_code == 200
            assert settings_response.json()["runtime_settings_writable"] is False

            response = client.put(
                "/api/settings/vlm",
                json={
                    "api_key": "test-secret",
                    "model_name": "test-model",
                    "libreoffice_path": "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                    "provider": "openai",
                },
            )
            assert response.status_code == 403
            assert not env_path.exists()
    finally:
        get_settings.cache_clear()


def test_cors_origin_env_parser() -> None:
    from app.config import parse_cors_allowed_origins, resolved_cors_allow_origin_regex

    assert parse_cors_allowed_origins("https://app.example.com, https://admin.example.com") == [
        "https://app.example.com",
        "https://admin.example.com",
    ]
    default_origins = parse_cors_allowed_origins("")
    assert "http://127.0.0.1:5173" in default_origins
    assert "http://localhost:4173" in default_origins
    default_regex = resolved_cors_allow_origin_regex(None)
    assert default_regex is not None
    assert re.match(default_regex, "http://127.0.0.1:53123")
    assert re.match(default_regex, "http://localhost:53123")
    assert not re.match(default_regex, "https://example.com")


def test_production_env_examples_disable_local_cors_regex() -> None:
    from app.config import resolved_cors_allow_origin_regex

    root = Path(__file__).resolve().parents[2]
    production_values = _read_env_example(root / "backend/.env.production.example")
    assert production_values["APP_ENV"] == "production"
    cors_regex = resolved_cors_allow_origin_regex(production_values["CORS_ALLOW_ORIGIN_REGEX"])
    assert cors_regex == "^$"
    assert not re.match(cors_regex, "http://localhost:53123")
    assert not re.match(cors_regex, "http://127.0.0.1:53123")

    local_values = _read_env_example(root / ".env.example")
    assert local_values["APP_ENV"] == "development"
    assert local_values["VLM_PROVIDER"] == "mock"
    assert json.loads(local_values["VLM_INFERENCE_PARAMS"].strip("'\""))["reasoning_effort"] == "off"
    assert json.loads(production_values["VLM_INFERENCE_PARAMS"].strip("'\""))["thinking"] == "off"


def test_production_security_headers(monkeypatch) -> None:
    from fastapi.testclient import TestClient

    from app.main import app

    monkeypatch.setenv("APP_ENV", "production")
    monkeypatch.setenv("SECURITY_HEADERS_ENABLED", "true")
    get_settings.cache_clear()

    try:
        with TestClient(app, base_url="https://testserver") as client:
            health = client.get("/api/health")
            assert health.status_code == 200
            assert health.headers["x-content-type-options"] == "nosniff"
            assert health.headers["x-frame-options"] == "DENY"
            assert health.headers["referrer-policy"] == "strict-origin-when-cross-origin"
            assert "camera=()" in health.headers["permissions-policy"]
            assert "default-src 'self'" in health.headers["content-security-policy"]
    finally:
        monkeypatch.delenv("APP_ENV", raising=False)
        monkeypatch.delenv("SECURITY_HEADERS_ENABLED", raising=False)
        get_settings.cache_clear()


def test_bank_poc_template_seed_is_idempotent() -> None:
    with get_client() as client:
        seeded = client.post("/api/templates/bank-documents-poc/seed")
        assert seeded.status_code == 200, seeded.text
        payload = seeded.json()
        assert payload["template_key"] == "bank_documents_poc"
        assert all(payload["created"].values())
        assert payload["schema"]["name"] == "은행 서류 핵심 정보"
        assert [field["key_name"] for field in payload["schema"]["fields"]] == ["문서번호", "고객명", "신청일", "금액", "발급기관"]
        assert payload["classifier"]["name"] == "bank_document_classifier"
        assert [item["class_name"] for item in payload["classifier"]["classes"]] == ["신청서", "동의서", "증빙문서"]
        assert payload["checklist"]["name"] == "bank_required_fields"
        assert payload["workflow"]["name"] == "은행 서류 자동 분류 및 검수"
        assert payload["workflow"]["validation_warnings"] == []
        definition = payload["workflow"]["definition"]
        nodes_by_id = {node["id"]: node for node in definition["nodes"]}
        assert nodes_by_id["classifier"]["data"]["config"]["classifier_id"] == payload["classifier"]["id"]
        assert nodes_by_id["branch"]["data"]["branchKeys"] == [
            "class:신청서",
            "class:동의서",
            "class:증빙문서",
            "unknown",
        ]
        assert nodes_by_id["kie_application"]["data"]["config"]["schema_id"] == payload["schema"]["id"]
        assert nodes_by_id["kie_supporting"]["data"]["config"]["schema_id"] == payload["schema"]["id"]
        assert nodes_by_id["required_application"]["data"]["config"]["checklist_id"] == payload["checklist"]["id"]
        assert nodes_by_id["required_consent"]["data"]["config"]["checklist_id"] == payload["checklist"]["id"]
        assert nodes_by_id["input"]["position"] == {"x": 40, "y": 240}
        assert nodes_by_id["classifier"]["position"] == {"x": 250, "y": 240}
        assert nodes_by_id["branch"]["position"] == {"x": 460, "y": 210}
        assert nodes_by_id["kie_application"]["position"] == {"x": 670, "y": 70}
        assert nodes_by_id["required_application"]["position"] == {"x": 900, "y": 70}
        assert nodes_by_id["required_consent"]["position"] == {"x": 670, "y": 240}
        assert nodes_by_id["kie_supporting"]["position"] == {"x": 670, "y": 410}
        assert nodes_by_id["merge"]["position"] == {"x": 900, "y": 300}
        assert nodes_by_id["export"]["position"] == {"x": 1110, "y": 300}
        assert nodes_by_id["kie_application"]["position"]["y"] < nodes_by_id["required_consent"]["position"]["y"] < nodes_by_id["kie_supporting"]["position"]["y"]
        assert {
            edge.get("sourceHandle")
            for edge in definition["edges"]
            if edge["source"] == "branch"
        } == {
            "class:신청서",
            "class:동의서",
            "class:증빙문서",
            "unknown",
        }
        sample_documents = payload["sample_documents"]
        assert len(sample_documents) == 3
        assert payload["sample_document"]["document_id"] == sample_documents[0]["document_id"]
        assert [document["library_path"] for document in sample_documents] == [
            "bank-poc/bank_00006.jpg",
            "bank-poc/bank_00008.jpg",
            "bank-poc/bank_00018.jpeg",
        ]
        assert [document["source_path"] for document in sample_documents] == [
            "assets/sample/bank_00006.jpg",
            "assets/sample/bank_00008.jpg",
            "assets/sample/bank_00018.jpeg",
        ]
        assert all(document["source_note"] == "사용자 제공 로컬 샘플 자산" for document in sample_documents)
        assert all(document["status"] == "ready" for document in sample_documents)

        reseeded = client.post("/api/templates/bank-documents-poc/seed")
        assert reseeded.status_code == 200, reseeded.text
        reseeded_payload = reseeded.json()
        assert not any(reseeded_payload["created"].values())
        assert reseeded_payload["schema"]["id"] == payload["schema"]["id"]
        assert reseeded_payload["workflow"]["id"] == payload["workflow"]["id"]
        assert [document["document_id"] for document in reseeded_payload["sample_documents"]] == [
            document["document_id"] for document in sample_documents
        ]


def test_bank_poc_seed_migrates_legacy_english_schema() -> None:
    with get_client() as client:
        legacy = client.post(
            "/api/schemas",
            json={
                "name": "bank_document_core_fields",
                "display_name": "은행 서류 핵심 정보",
                "fields": [
                    {"key_name": "document_number", "description": "legacy", "output_format": "string"},
                ],
            },
        )
        assert legacy.status_code == 200, legacy.text
        legacy_classifier = client.post(
            "/api/document-classifiers",
            json={
                "name": "bank_document_classifier",
                "description": "legacy english classifier",
                "allow_unknown": True,
                "classes": [
                    {"class_name": "application_form", "description": "legacy", "signals": ["application"]},
                    {"class_name": "consent_form", "description": "legacy", "signals": ["consent"]},
                    {"class_name": "supporting_document", "description": "legacy", "signals": ["support"]},
                ],
            },
        )
        assert legacy_classifier.status_code == 200, legacy_classifier.text
        legacy_checklist = client.post(
            "/api/required-field-checklists",
            json={
                "name": "bank_required_fields",
                "description": "legacy checklist",
                "items": [
                    {
                        "item_name": "legacy_item",
                        "description": "legacy",
                        "evidence_type": "text_or_handwriting",
                        "required": True,
                    },
                ],
            },
        )
        assert legacy_checklist.status_code == 200, legacy_checklist.text

        seeded = client.post("/api/templates/bank-documents-poc/seed")
        assert seeded.status_code == 200, seeded.text
        payload = seeded.json()
        assert payload["created"]["schema"] is False
        assert payload["created"]["classifier"] is False
        assert payload["created"]["checklist"] is False
        assert payload["schema"]["name"] == "은행 서류 핵심 정보"
        assert [field["key_name"] for field in payload["schema"]["fields"]] == ["문서번호", "고객명", "신청일", "금액", "발급기관"]
        assert [item["class_name"] for item in payload["classifier"]["classes"]] == ["신청서", "동의서", "증빙문서"]
        assert [item["item_name"] for item in payload["checklist"]["items"]] == ["고객명", "작성일", "필수 동의 체크", "서명 또는 날인"]
        nodes_by_id = {node["id"]: node for node in payload["workflow"]["definition"]["nodes"]}
        assert nodes_by_id["branch"]["data"]["branchKeys"] == ["class:신청서", "class:동의서", "class:증빙문서", "unknown"]
        assert nodes_by_id["kie_application"]["data"]["config"]["schema_id"] == payload["schema"]["id"]
        assert nodes_by_id["kie_supporting"]["data"]["config"]["schema_id"] == payload["schema"]["id"]
        assert {
            edge.get("sourceHandle")
            for edge in payload["workflow"]["definition"]["edges"]
            if edge["source"] == "branch"
        } == {"class:신청서", "class:동의서", "class:증빙문서", "unknown"}
        active_schema_names = [schema["name"] for schema in client.get("/api/schemas").json()]
        assert "은행 서류 핵심 정보" in active_schema_names
        assert "bank_document_core_fields" not in active_schema_names


def test_upload_rejects_content_type_spoof() -> None:
    with get_client() as client:
        response = client.post(
            "/api/documents",
            files={"file": ("fake.png", b"not an image", "image/png")},
        )
        assert response.status_code == 415


def test_upload_rejects_file_size_over_limit(monkeypatch) -> None:
    monkeypatch.setenv("UPLOAD_MAX_FILE_BYTES", "4")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            response = client.post(
                "/api/documents",
                files={"file": ("invoice.png", ONE_BY_ONE_PNG, "image/png")},
            )
            assert response.status_code == 413
    finally:
        monkeypatch.delenv("UPLOAD_MAX_FILE_BYTES", raising=False)
        get_settings.cache_clear()


def test_upload_downscales_preview_over_pixel_limit(monkeypatch) -> None:
    image = Image.new("RGB", (20, 10), (255, 255, 255))
    stream = io.BytesIO()
    image.save(stream, format="PNG")

    monkeypatch.setenv("UPLOAD_MAX_IMAGE_PIXELS", "100")
    monkeypatch.setenv("DOCUMENT_PAGE_MAX_LONG_EDGE", "0")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            response = client.post(
                "/api/documents",
                files={"file": ("oversized.png", stream.getvalue(), "image/png")},
            )
            assert response.status_code == 200, response.text
            page = response.json()["pages"][0]
            assert page["width"] * page["height"] <= 100
            assert page["width"] * page["height"] <= 98
    finally:
        monkeypatch.delenv("UPLOAD_MAX_IMAGE_PIXELS", raising=False)
        monkeypatch.delenv("DOCUMENT_PAGE_MAX_LONG_EDGE", raising=False)
        get_settings.cache_clear()


def test_batch_upload_rejects_too_many_files(monkeypatch) -> None:
    monkeypatch.setenv("UPLOAD_MAX_BATCH_FILES", "1")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            schema = create_schema(client, "batch_limit_schema")
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", ("a.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("b.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 413
    finally:
        monkeypatch.delenv("UPLOAD_MAX_BATCH_FILES", raising=False)
        get_settings.cache_clear()


def test_large_batch_multipart_allows_configured_counts_above_parser_default(monkeypatch) -> None:
    monkeypatch.setenv("UPLOAD_MAX_BATCH_FILES", "10000")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            files = [
                ("files", (f"document_{index:04d}.png", ONE_BY_ONE_PNG, "image/png"))
                for index in range(1001)
            ]
            response = client.post(
                "/api/classification-batches",
                data={"classifier_id": "missing"},
                files=files,
            )
            assert response.status_code == 404, response.text
            assert response.json()["detail"] == "Document classifier not found"
    finally:
        monkeypatch.delenv("UPLOAD_MAX_BATCH_FILES", raising=False)
        get_settings.cache_clear()


def test_large_batch_multipart_returns_413_over_configured_limit(monkeypatch) -> None:
    monkeypatch.setenv("UPLOAD_MAX_BATCH_FILES", "1001")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            files = [
                ("files", (f"document_{index:04d}.png", ONE_BY_ONE_PNG, "image/png"))
                for index in range(1002)
            ]
            response = client.post(
                "/api/classification-batches",
                data={"classifier_id": "missing"},
                files=files,
            )
            assert response.status_code == 413, response.text
            assert "1001" in response.json()["detail"]
    finally:
        monkeypatch.delenv("UPLOAD_MAX_BATCH_FILES", raising=False)
        get_settings.cache_clear()


def test_retention_cleanup_removes_expired_uploads(monkeypatch) -> None:
    from app.database import SessionLocal

    monkeypatch.setenv("UPLOAD_RETENTION_HOURS", "1")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            document = upload_png(client)
            db = SessionLocal()
            try:
                row = db.get(Document, document["document_id"])
                assert row is not None
                storage_path = row.storage_path
                extra_root = Path(storage_path).parent.parent / "retention-extra"
                raw_root = extra_root / "raw-expired"
                export_root = extra_root / "exports"
                raw_root.mkdir(parents=True, exist_ok=True)
                export_root.mkdir(parents=True, exist_ok=True)
                raw_source = raw_root / "content.bin"
                raw_pdf = raw_root / "preview.pdf"
                raw_html = raw_root / "preview.html"
                export_file = export_root / "workflow.xlsx"
                for path in [raw_source, raw_pdf, raw_html, export_file]:
                    path.write_bytes(b"expired")
                row.created_at = datetime.utcnow() - timedelta(hours=2)
                for page in row.pages:
                    page.created_at = datetime.utcnow() - timedelta(hours=2)
                raw = RawExtraction(
                    filename="expired-source.docx",
                    source_format="docx",
                    size_bytes=7,
                    storage_path=str(raw_source),
                    pdf_path=str(raw_pdf),
                    html_path=str(raw_html),
                    status="completed",
                    created_at=datetime.utcnow() - timedelta(hours=2),
                )
                export = ExportJob(
                    owner_type="workflow_run",
                    owner_id="expired-run",
                    format="xlsx",
                    status="completed",
                    filename="workflow.xlsx",
                    storage_path=str(export_file),
                    content_type=XLSX_MIME,
                    size_bytes=7,
                    created_at=datetime.utcnow() - timedelta(hours=2),
                )
                audit = AuditEvent(
                    entity_type="workflow_run",
                    entity_id="expired-run",
                    action="completed",
                    message="expired audit",
                    created_at=datetime.utcnow() - timedelta(hours=2),
                )
                db.add_all([raw, export, audit])
                db.commit()
                raw_id = raw.id
                export_id = export.id
                audit_id = audit.id
            finally:
                db.close()

            cleaned = client.post("/api/maintenance/retention-cleanup")
            assert cleaned.status_code == 200, cleaned.text
            payload = cleaned.json()
            assert payload["status"] == "cleaned"
            assert payload["counts"]["documents"] >= 1
            assert payload["counts"]["raw_extractions"] >= 1
            assert payload["counts"]["export_jobs"] >= 1
            assert payload["counts"]["audit_events"] >= 1
            assert client.get(f"/api/documents/{document['document_id']}").status_code == 404
            assert not os.path.exists(os.path.dirname(storage_path))
            assert not raw_root.exists()
            assert not export_root.exists()
            db = SessionLocal()
            try:
                assert db.get(RawExtraction, raw_id) is None
                assert db.get(ExportJob, export_id) is None
                assert db.get(AuditEvent, audit_id) is None
            finally:
                db.close()
    finally:
        monkeypatch.delenv("UPLOAD_RETENTION_HOURS", raising=False)
        get_settings.cache_clear()


def test_vlm_runtime_kwargs_include_speed_controls(monkeypatch) -> None:
    from app.vlm import _build_llm_kwargs

    try:
        monkeypatch.setenv("VLM_API_KEY", "test-secret")
        monkeypatch.setenv("VLM_MODEL_NAME", "test-model")
        monkeypatch.setenv(
            "VLM_INFERENCE_PARAMS",
            json.dumps(
                {
                    "reasoning_effort": "minimal",
                    "thinking": "minimal",
                    "temperature": "0.2",
                    "verbosity": "low",
                    "max_completion_tokens": "1024",
                    "top_p": "0.8",
                    "service_tier": "auto",
                }
            ),
        )
        get_settings.cache_clear()

        kwargs = _build_llm_kwargs()
        assert kwargs["reasoning_effort"] == "minimal"
        assert kwargs["temperature"] == 0.2
        assert kwargs["verbosity"] == "low"
        assert kwargs["max_completion_tokens"] == 1024
        assert kwargs["top_p"] == 0.8
        assert kwargs["service_tier"] == "auto"
    finally:
        get_settings.cache_clear()


def test_vlm_inference_params_default_reasoning_off(monkeypatch) -> None:
    from app.vlm import _build_google_generation_config, _build_llm_kwargs

    try:
        for key in [
            "VLM_INFERENCE_PARAMS",
            "VLM_REASONING_EFFORT",
            "VLM_VERBOSITY",
            "VLM_MAX_COMPLETION_TOKENS",
            "VLM_TOP_P",
            "VLM_SERVICE_TIER",
        ]:
            monkeypatch.delenv(key, raising=False)
        monkeypatch.setenv("VLM_API_KEY", "test-secret")
        monkeypatch.setenv("VLM_MODEL_NAME", "test-model")
        get_settings.cache_clear()

        kwargs = _build_llm_kwargs()
        assert "reasoning_effort" not in kwargs
        assert "verbosity" not in kwargs
        assert kwargs["temperature"] == 0

        schema = {"type": "object", "properties": {"name": {"type": "string"}}, "required": ["name"]}
        config = _build_google_generation_config("system", schema)
        assert config["thinking_config"] == {"thinking_budget": 0}
    finally:
        get_settings.cache_clear()


def test_kie_field_groups_run_with_bounded_parallelism(monkeypatch) -> None:
    from app import extraction as extraction_module
    from app.extraction import DocumentPageSnapshot, DocumentSnapshot
    from app.schemas import FieldDefinition

    monkeypatch.setenv("KIE_FIELD_GROUP_SIZE", "1")
    monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "1")
    monkeypatch.setenv("VLM_MAX_CONCURRENT_REQUESTS", "2")
    get_settings.cache_clear()

    active = 0
    max_active = 0
    calls = 0
    lock = threading.Lock()

    def fake_extract_with_vlm(fields, image_inputs=None):
        nonlocal active, max_active, calls
        with lock:
            active += 1
            calls += 1
            max_active = max(max_active, active)
        time.sleep(0.03)
        with lock:
            active -= 1
        return {field.key_name: {"value": f"value-{field.key_name}", "page": 1, "confidence": 1, "evidence": "mock"} for field in fields}

    try:
        monkeypatch.setattr(extraction_module, "extract_with_vlm", fake_extract_with_vlm)
        fields = [
            FieldDefinition(key_name=f"field_{index}", description="test field", output_format="string")
            for index in range(4)
        ]
        document = DocumentSnapshot(
            id="doc_1",
            storage_path="/tmp/doc_1/source.png",
            pages=[DocumentPageSnapshot(page_number=1, image_path="/tmp/doc_1/page_1.png")],
        )

        values = extraction_module._extract_grouped_values(document, fields, [], "job_1")

        assert calls == 4
        assert max_active == 2
        assert set(values) == {field.key_name for field in fields}
    finally:
        get_settings.cache_clear()


def test_module_batch_workers_use_workflow_worker_limit(monkeypatch) -> None:
    from app import document_modules

    active = 0
    max_active = 0
    finalized = False
    lock = threading.Lock()

    def fake_runner(_job_id: str) -> None:
        nonlocal active, max_active
        with lock:
            active += 1
            max_active = max(max_active, active)
        time.sleep(0.03)
        with lock:
            active -= 1

    def fake_failer(_job_id: str, message: str) -> None:
        raise AssertionError(message)

    def fake_finalizer() -> None:
        nonlocal finalized
        finalized = True

    monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
    get_settings.cache_clear()
    try:
        document_modules._run_parallel_batch(["job_1", "job_2", "job_3", "job_4"], fake_runner, fake_failer, fake_finalizer)
        assert max_active == 2
        assert finalized is True
    finally:
        get_settings.cache_clear()


def test_module_batch_async_workers_use_workflow_worker_limit(monkeypatch) -> None:
    from app import document_modules

    active = 0
    max_active = 0
    finalized = False

    async def fake_runner(_job_id: str) -> None:
        nonlocal active, max_active
        active += 1
        max_active = max(max_active, active)
        await asyncio.sleep(0.01)
        active -= 1

    def fake_failer(_job_id: str, message: str) -> None:
        raise AssertionError(message)

    def fake_finalizer() -> None:
        nonlocal finalized
        finalized = True

    monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
    get_settings.cache_clear()
    try:
        asyncio.run(document_modules._run_parallel_batch_async(["job_1", "job_2", "job_3", "job_4"], fake_runner, fake_failer, fake_finalizer))
        assert max_active == 2
        assert finalized is True
    finally:
        get_settings.cache_clear()


def test_extraction_batch_workers_use_workflow_worker_limit(monkeypatch) -> None:
    from app import extraction as extraction_module

    active = 0
    max_active = 0
    finalized: list[str] = []
    lock = threading.Lock()

    def fake_run_job(_job_id: str) -> None:
        nonlocal active, max_active
        with lock:
            active += 1
            max_active = max(max_active, active)
        time.sleep(0.03)
        with lock:
            active -= 1

    monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
    monkeypatch.setattr(extraction_module, "run_extraction_job", fake_run_job)
    monkeypatch.setattr(extraction_module, "_finalize_batch", lambda batch_id: finalized.append(batch_id))
    get_settings.cache_clear()
    try:
        extraction_module.run_batch_jobs("batch_1", ["job_1", "job_2", "job_3", "job_4"])
        assert max_active == 2
        assert finalized == ["batch_1"]
    finally:
        get_settings.cache_clear()


def test_extraction_batch_async_workers_use_workflow_worker_limit(monkeypatch) -> None:
    from app import extraction as extraction_module

    active = 0
    max_active = 0
    finalized: list[str] = []

    async def fake_run_job(_job_id: str) -> None:
        nonlocal active, max_active
        active += 1
        max_active = max(max_active, active)
        await asyncio.sleep(0.01)
        active -= 1

    monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
    monkeypatch.setattr(extraction_module, "_run_batch_extraction_job_async", fake_run_job)
    monkeypatch.setattr(extraction_module, "_finalize_batch", lambda batch_id: finalized.append(batch_id))
    get_settings.cache_clear()
    try:
        asyncio.run(extraction_module.run_batch_jobs_async("batch_1", ["job_1", "job_2", "job_3", "job_4"]))
        assert max_active == 2
        assert finalized == ["batch_1"]
    finally:
        get_settings.cache_clear()


def test_vlm_api_style_auto_detects_google_and_base_url(monkeypatch) -> None:
    from app.vlm import resolve_vlm_api_style

    try:
        monkeypatch.setenv("VLM_PROVIDER", "auto")
        monkeypatch.setenv("VLM_API_KEY", "AIzaSyCP_test_key")
        monkeypatch.setenv("VLM_MODEL_NAME", "gemini-3.1-flash-lite")
        monkeypatch.setenv("VLM_BASE_URL", "")
        get_settings.cache_clear()
        assert resolve_vlm_api_style() == "google_genai"

        monkeypatch.setenv("VLM_BASE_URL", "https://openrouter.ai/api/v1")
        get_settings.cache_clear()
        assert resolve_vlm_api_style() == "openai_compatible"

        monkeypatch.setenv("VLM_PROVIDER", "openai")
        monkeypatch.setenv("VLM_BASE_URL", "")
        get_settings.cache_clear()
        assert resolve_vlm_api_style() == "google_genai"
    finally:
        get_settings.cache_clear()


def test_vlm_base_url_allows_local_provider_without_api_key(monkeypatch) -> None:
    from app.vlm import _build_llm_kwargs, _ensure_vlm_credentials

    try:
        monkeypatch.setenv("VLM_PROVIDER", "auto")
        monkeypatch.setenv("VLM_API_KEY", "")
        monkeypatch.setenv("OPENAI_API_KEY", "")
        monkeypatch.setenv("VLM_MODEL_NAME", "local-model")
        monkeypatch.setenv("VLM_BASE_URL", "http://127.0.0.1:11434/v1")
        get_settings.cache_clear()

        settings = get_settings()
        _ensure_vlm_credentials(settings)
        kwargs = _build_llm_kwargs()
        assert kwargs["base_url"] == "http://127.0.0.1:11434/v1"
        assert kwargs["api_key"] == "local-vlm"
    finally:
        get_settings.cache_clear()


def test_vlm_base_url_does_not_replace_google_api_key(monkeypatch) -> None:
    from app.vlm import VlmRuntimeError, _ensure_vlm_credentials

    try:
        monkeypatch.setenv("VLM_PROVIDER", "google_genai")
        monkeypatch.setenv("VLM_API_KEY", "")
        monkeypatch.setenv("GOOGLE_API_KEY", "")
        monkeypatch.setenv("VLM_MODEL_NAME", "gemini-test-model")
        monkeypatch.setenv("VLM_BASE_URL", "http://127.0.0.1:11434/v1")
        get_settings.cache_clear()

        settings = get_settings()
        with pytest.raises(VlmRuntimeError) as exc:
            _ensure_vlm_credentials(settings)
        assert exc.value.code == "VLM_CREDENTIALS_MISSING"
        assert "VLM API key is required" in str(exc.value)
    finally:
        get_settings.cache_clear()


def test_vlm_errors_have_stable_codes_and_redact_secrets(monkeypatch) -> None:
    from app.vlm import VlmRuntimeError, _coerce_structured_response, _sanitize_provider_error, resolve_vlm_api_style

    try:
        monkeypatch.setenv("VLM_PROVIDER", "unknown_provider")
        get_settings.cache_clear()
        with pytest.raises(VlmRuntimeError) as provider_error:
            resolve_vlm_api_style()
        assert provider_error.value.code == "VLM_PROVIDER_UNSUPPORTED"
        assert provider_error.value.as_detail()["code"] == "VLM_PROVIDER_UNSUPPORTED"

        monkeypatch.setenv("VLM_API_KEY", "AIzaSyCP_test_key_should_not_leak")
        get_settings.cache_clear()
        sanitized = _sanitize_provider_error(RuntimeError("bad key AIzaSyCP_test_key_should_not_leak"))
        assert "AIzaSyCP_test_key_should_not_leak" not in sanitized
        assert "[redacted]" in sanitized

        with pytest.raises(VlmRuntimeError) as response_error:
            _coerce_structured_response(SimpleNamespace(text="not-json"))
        assert response_error.value.code == "VLM_RESPONSE_INVALID_JSON"
    finally:
        get_settings.cache_clear()


def test_vlm_provider_request_retries_transient_broken_pipe(monkeypatch) -> None:
    from app import vlm as vlm_module
    from app.vlm import VlmRuntimeError, _invoke_vlm_with_limit

    calls = 0

    def fake_invoke(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        if calls == 1:
            raise VlmRuntimeError("VLM_PROVIDER_REQUEST_FAILED", "Google GenAI VLM request failed: [Errno 32] Broken pipe")
        return {"ok": True}

    try:
        monkeypatch.setenv("VLM_MAX_RETRIES", "2")
        get_settings.cache_clear()
        monkeypatch.setattr(vlm_module.time, "sleep", lambda _delay: None)
        monkeypatch.setattr(vlm_module, "_invoke_structured_llm", fake_invoke)

        assert _invoke_vlm_with_limit("system", "prompt", [], {}, "google_genai") == {"ok": True}
        assert calls == 2
    finally:
        get_settings.cache_clear()


def test_google_generation_config_uses_structured_output_and_explicit_thinking_level(monkeypatch) -> None:
    from app.vlm import _build_google_generation_config

    try:
        monkeypatch.setenv("VLM_PROVIDER", "auto")
        monkeypatch.setenv("VLM_API_KEY", "AIzaSyCP_test_key")
        monkeypatch.setenv("VLM_MODEL_NAME", "gemini-3.1-flash-lite")
        monkeypatch.setenv("VLM_INFERENCE_PARAMS", json.dumps({"reasoning_effort": "minimal", "thinking": "minimal"}))
        get_settings.cache_clear()

        schema = {"type": "object", "properties": {"name": {"type": "string"}}, "required": ["name"]}
        config = _build_google_generation_config("system", schema)
        assert config["system_instruction"] == "system"
        assert config["response_mime_type"] == "application/json"
        assert config["response_json_schema"] == schema
        assert config["thinking_config"] == {"thinking_level": "minimal"}
    finally:
        get_settings.cache_clear()


def test_classification_schema_uses_classes_or_unknown_only() -> None:
    from app.vlm import _classification_output_schema
    from app.schemas import ClassCandidate

    spec = _classification_output_schema(
        classes=[ClassCandidate(class_name="contract", description="Contract", signals=[])],
        allow_unknown=False,
    )
    schema = spec.json_schema()
    assert schema["properties"]["status"]["enum"] == ["classified", "unknown"]
    class_name_schema = schema["properties"]["class_name"]
    assert {"type": "null"} in class_name_schema["anyOf"]


def test_dynamic_pydantic_output_schemas_require_configured_keys() -> None:
    from app.vlm import _required_field_output_schema, build_structured_output_schema
    from app.schemas import FieldDefinition, RequiredFieldItem

    kie_spec = build_structured_output_schema(
        [
            FieldDefinition(key_name="본인 성명", description="우측 하단 본인 성명", output_format="string"),
            FieldDefinition(key_name="금액", description="금액", output_format="float"),
        ]
    )
    kie_schema = kie_spec.json_schema()
    assert kie_schema["required"] == ["본인 성명", "금액"]
    assert set(kie_schema["properties"]) == {"본인 성명", "금액"}

    required_spec = _required_field_output_schema(
        [
            RequiredFieldItem(item_name="서명", description="서명 확인"),
            RequiredFieldItem(item_name="작성일", description="작성일 확인"),
        ]
    )
    required_schema = required_spec.json_schema()
    item_ref = required_schema["properties"]["items"]["$ref"]
    ref_name = item_ref.removeprefix("#/$defs/")
    item_schema = required_schema["$defs"][ref_name]
    assert item_schema["required"] == ["서명", "작성일"]
    assert set(item_schema["properties"]) == {"서명", "작성일"}


def test_classification_validation_clears_unknown_class_name() -> None:
    from app.document_modules import ClassificationContext, _validate_classification_output
    from app.extraction import DocumentSnapshot
    from app.schemas import ClassCandidate

    context = ClassificationContext(
        document=DocumentSnapshot(id="doc_1", storage_path="", pages=[]),
        classifier_id="clf_1",
        classes=[ClassCandidate(class_name="contract", description="Contract", signals=[])],
        allow_unknown=True,
    )
    output = _validate_classification_output(
        {
            "status": "unknown",
            "class_name": "contract",
            "confidence": 0.3,
            "reason": "No class matched.",
            "evidence": [],
        },
        context,
    )
    assert output["status"] == "unknown"
    assert output["class_name"] is None


def test_schema_validation_and_creation() -> None:
    with get_client() as client:
        invalid = client.post(
            "/api/schemas",
            json={
                "name": "bad_schema",
                "fields": [
                    {"key_name": "total", "description": "Total amount", "output_format": "float"},
                    {"key_name": "total", "description": "Duplicate total", "output_format": "float"},
                ],
            },
        )
        assert invalid.status_code == 422

        unsupported_format = client.post(
            "/api/schemas",
            json={
                "name": "bad_format",
                "fields": [
                    {"key_name": "count", "description": "Unsupported integer field", "output_format": "int"},
                ],
            },
        )
        assert unsupported_format.status_code == 422

        valid = create_schema(client)
        assert valid["name"] == "invoice_basic"
        assert valid["fields"][0]["key_name"] == "invoice_number"

        region_schema = client.post(
            "/api/schemas",
            json={
                "name": "region_schema",
                "regions": [
                    {"id": "region_1", "name": "Region 1", "page": 1, "x": 0.1, "y": 0.2, "width": 0.3, "height": 0.1}
                ],
                "fields": [
                    {
                        "key_name": "handwritten_name",
                        "description": "손글씨 이름 영역",
                        "output_format": "string",
                        "region_id": "region_1",
                    },
                ],
            },
        )
        assert region_schema.status_code == 200, region_schema.text
        assert region_schema.json()["regions"][0]["x"] == 0.1
        assert region_schema.json()["fields"][0]["region_id"] == "region_1"

        invalid_region = client.post(
            "/api/schemas",
            json={
                "name": "invalid_region_schema",
                "regions": [
                    {"id": "region_1", "name": "Region 1", "page": 1, "x": 0.8, "y": 0.2, "width": 0.3, "height": 0.1}
                ],
                "fields": [
                    {
                        "key_name": "handwritten_name",
                        "description": "손글씨 이름 영역",
                        "output_format": "string",
                        "region_id": "region_1",
                    },
                ],
            },
        )
        assert invalid_region.status_code == 422

        korean_with_space = client.post(
            "/api/schemas",
            json={
                "name": "korean_schema",
                "fields": [
                    {"key_name": "법정 대리인 성", "description": "우측 하단의 법정 대리인 성명", "output_format": "string"},
                ],
            },
        )
        assert korean_with_space.status_code == 200
        assert korean_with_space.json()["fields"][0]["key_name"] == "법정 대리인 성"

        screenshot_payload = client.post(
            "/api/schemas",
            json={
                "name": "document_schema",
                "fields": [
                    {"key_name": "개정일", "description": "좌측 하단의 개정일자", "output_format": "date"},
                    {"key_name": "본인 성명", "description": "우측 하단의 본인 성명", "output_format": "string"},
                ],
            },
        )
        assert screenshot_payload.status_code == 200, screenshot_payload.text
        assert screenshot_payload.json()["fields"][1]["key_name"] == "본인 성명"


def test_image_upload() -> None:
    with get_client() as client:
        document = upload_png(client)
        assert document["page_count"] == 1
        assert document["created_at"]
        image = client.get(document["pages"][0]["image_url"])
        assert image.status_code == 200
        assert image.headers["content-type"] == "image/jpeg"
        thumbnail = client.get(f"/api/documents/{document['document_id']}/pages/1/thumbnail?width=96")
        assert thumbnail.status_code == 200
        assert thumbnail.headers["content-type"] == "image/jpeg"

        documents = client.get("/api/documents").json()
        assert any(item["document_id"] == document["document_id"] for item in documents)


def test_document_library_select_copy_move_and_folder_operations() -> None:
    from app.database import SessionLocal

    with get_client() as client:
        document = upload_png(client)
        document_id = document["document_id"]

        created_folder = client.post("/api/library/folders", json={"folder_path": "검수"})
        assert created_folder.status_code == 200, created_folder.text
        assert any(folder["path"] == "검수" for folder in created_folder.json()["folders"])

        moved = client.post(
            "/api/library/move",
            json={"document_ids": [document_id], "target_folder": "검수"},
        )
        assert moved.status_code == 200, moved.text
        moved_document = moved.json()["documents"][0]
        assert moved_document["document_id"] == document_id
        assert moved_document["library_path"] == "검수/invoice.png"

        ids = client.get("/api/documents/ids?library_path=검수").json()
        assert document_id in ids

        db = SessionLocal()
        try:
            original_row = db.get(Document, document_id)
            assert original_row is not None
            original_storage_path = original_row.storage_path
        finally:
            db.close()

        copied = client.post(
            "/api/library/copy",
            json={"document_ids": [document_id], "target_folder": "검수"},
        )
        assert copied.status_code == 200, copied.text
        copied_document = copied.json()["documents"][0]
        copied_document_id = copied_document["document_id"]
        assert copied_document_id != document_id
        assert copied_document["library_path"].startswith("검수/invoice copy")

        db = SessionLocal()
        try:
            copied_row = db.get(Document, copied_document_id)
            assert copied_row is not None
            assert copied_row.storage_path != original_storage_path
            assert copied_row.page_count == 1
            assert len(copied_row.pages) == 1
        finally:
            db.close()

        folder_copy = client.post(
            "/api/library/copy",
            json={"folder_paths": ["검수"], "target_folder": ""},
        )
        assert folder_copy.status_code == 200, folder_copy.text
        assert any(item["library_path"].startswith("검수 copy/") for item in folder_copy.json()["documents"])

        moved_folder = client.post(
            "/api/library/move",
            json={"folder_paths": ["검수 copy"], "target_folder": "완료"},
        )
        assert moved_folder.status_code == 200, moved_folder.text
        assert any(item["library_path"].startswith("완료/검수 copy/") for item in moved_folder.json()["documents"])

        deleted_original = client.delete(f"/api/documents/{document_id}")
        assert deleted_original.status_code == 200, deleted_original.text
        copied_payload = client.get(f"/api/documents/{copied_document_id}")
        assert copied_payload.status_code == 200, copied_payload.text
        page_response = client.get(copied_payload.json()["pages"][0]["image_url"])
        assert page_response.status_code == 200

        bulk_deleted = client.post("/api/documents/delete", json={"document_ids": [copied_document_id]})
        assert bulk_deleted.status_code == 200, bulk_deleted.text
        assert bulk_deleted.json()["documents"][0]["status"] == "deleted"
        assert client.get(copied_payload.json()["pages"][0]["image_url"]).status_code == 410


def test_jpeg_upload_preserves_source_pixels_with_dpi_metadata() -> None:
    image = Image.new("RGB", (300, 420), (255, 255, 255))
    buffer = io.BytesIO()
    image.save(buffer, format="JPEG", dpi=(300, 300))

    with get_client() as client:
        response = client.post(
            "/api/documents",
            files={"file": ("scan.jpg", buffer.getvalue(), "image/jpeg")},
        )
        assert response.status_code == 200, response.text
        document = response.json()
        assert document["pages"][0]["width"] == 300
        assert document["pages"][0]["height"] == 420
        image_response = client.get(document["pages"][0]["image_url"])
        assert image_response.status_code == 200
        loaded = Image.open(io.BytesIO(image_response.content))
        assert loaded.size == (300, 420)


def test_pdf_upload() -> None:
    with get_client() as client:
        pdf_bytes = make_pdf_bytes()
        response = client.post(
            "/api/documents",
            files={"file": ("invoice.pdf", pdf_bytes, "application/pdf")},
        )
        assert response.status_code == 200
        document = response.json()
        assert document["page_count"] == 1
        assert document["pages"][0]["width"] > 0
        assert document["pages"][0]["height"] > 0


def test_office_upload_for_key_information_extractor(monkeypatch) -> None:
    def fake_convert(source_path, suffix, pdf_path):
        document = fitz.open()
        page = document.new_page(width=240, height=120)
        page.insert_text((24, 60), f"Converted {source_path.name}")
        document.save(pdf_path)
        document.close()

    monkeypatch.setattr("app.document_processor.convert_office_to_pdf", fake_convert)

    samples = [
        ("report.docx", make_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
        ("deck.pptx", make_pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    ]
    with get_client() as client:
        for filename, data, mime_type in samples:
            response = client.post(
                "/api/documents",
                files={"file": (filename, data, mime_type)},
            )
            assert response.status_code == 200, response.text
            document = response.json()
            assert document["filename"] == filename
            assert document["page_count"] == 1
            image = client.get(document["pages"][0]["image_url"])
            assert image.status_code == 200
            assert image.headers["content-type"] == "image/jpeg"


def test_extraction_fails_without_vlm_credentials(monkeypatch) -> None:
    monkeypatch.setenv("VLM_API_KEY", "")
    monkeypatch.setenv("VLM_MODEL_NAME", "")
    monkeypatch.setenv("OPENAI_API_KEY", "")
    monkeypatch.setenv("OPENAI_MODEL_NAME", "")
    get_settings.cache_clear()
    with get_client() as client:
        document = upload_png(client)
        schema = create_schema(client)
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema["id"]},
        )
        assert job_response.status_code == 200
        job_id = job_response.json()["job_id"]

        job = client.get(f"/api/extraction-jobs/{job_id}").json()
        assert job["status"] == "failed"
        assert job["result_id"] is None
        assert "VLM model name is required" in job["error_message"]


def test_schema_update_replaces_current_schema() -> None:
    with get_client() as client:
        schema = create_schema(client)
        updated = client.patch(
            f"/api/schemas/{schema['id']}",
            json={
                "display_name": "Updated Invoice Basic",
                "fields": [
                    {
                        "key_name": "invoice_number",
                        "description": "Invoice number near the top.",
                        "output_format": "string",
                    },
                    {
                        "key_name": "invoice_date",
                        "description": "Invoice issue date.",
                        "output_format": "date",
                    },
                ],
            },
        )
        assert updated.status_code == 200, updated.text
        payload = updated.json()
        assert "current_version" not in payload
        assert payload["display_name"] == "Updated Invoice Basic"
        assert payload["fields"][1]["key_name"] == "invoice_date"


def test_schema_update_allows_same_name_for_loaded_schema() -> None:
    with get_client() as client:
        schema = create_schema(client, name="테스트")
        updated = client.patch(
            f"/api/schemas/{schema['id']}",
            json={
                "name": "테스트",
                "display_name": "테스트",
                "description": "수정된 설명",
                "fields": [
                    {
                        "key_name": "수정필드",
                        "description": "사용자가 저장된 스키마를 불러온 뒤 수정한 필드",
                        "output_format": "string",
                    }
                ],
            },
        )
        assert updated.status_code == 200, updated.text
        payload = updated.json()
        assert payload["id"] == schema["id"]
        assert payload["name"] == "테스트"
        assert payload["description"] == "수정된 설명"
        assert "current_version" not in payload
        assert payload["fields"][0]["key_name"] == "수정필드"


def test_schema_update_merges_duplicate_loaded_schema_name() -> None:
    from app.database import SessionLocal
    from app.models import Schema

    with get_client() as client:
        schema = create_schema(client, name="중복스키마")
        duplicate_schema_json = {
            "name": "중복스키마",
            "display_name": "중복스키마",
            "description": "old duplicate",
            "is_template": False,
            "template_category": None,
            "pinned": False,
            "regions": [],
            "fields": [
                {
                    "key_name": "old_field",
                    "description": "Old duplicate field.",
                    "output_format": "string",
                }
            ],
        }
        db = SessionLocal()
        try:
            duplicate = Schema(
                name="중복스키마",
                display_name="중복스키마",
                description="old duplicate",
                current_version=1,
                schema_json=json.dumps(duplicate_schema_json, ensure_ascii=False),
                is_template=False,
                template_category=None,
                pinned=False,
                ephemeral=False,
            )
            db.add(duplicate)
            db.commit()
            duplicate_id = duplicate.id
        finally:
            db.close()

        updated = client.patch(
            f"/api/schemas/{schema['id']}",
            json={
                "name": "중복스키마",
                "display_name": "중복스키마",
                "description": "merged current schema",
                "fields": [
                    {
                        "key_name": "current_field",
                        "description": "Current schema field.",
                        "output_format": "string",
                    }
                ],
            },
        )
        assert updated.status_code == 200, updated.text
        schemas = [item for item in client.get("/api/schemas").json() if item["name"] == "중복스키마"]
        assert len(schemas) == 1
        assert schemas[0]["id"] == schema["id"]
        assert client.get(f"/api/schemas/{duplicate_id}").status_code == 404


def test_schema_delete_archives_and_allows_name_reuse() -> None:
    with get_client() as client:
        schema = create_schema(client, name="삭제테스트")
        deleted = client.delete(f"/api/schemas/{schema['id']}")
        assert deleted.status_code == 200, deleted.text
        assert deleted.json()["archived"] is True

        listed_names = [item["name"] for item in client.get("/api/schemas").json()]
        assert "삭제테스트" not in listed_names

        archived = client.get(f"/api/schemas/{schema['id']}")
        assert archived.status_code == 200
        assert archived.json()["archived"] is True

        recreated = create_schema(client, name="삭제테스트")
        assert recreated["id"] != schema["id"]


def test_schema_duplicate_copies_definition_with_incremented_name() -> None:
    with get_client() as client:
        schema = create_schema(client, name="성명")
        duplicated = client.post(f"/api/schemas/{schema['id']}/duplicate")
        assert duplicated.status_code == 200, duplicated.text
        payload = duplicated.json()
        assert payload["id"] != schema["id"]
        assert payload["name"] == "성명 (1)"
        assert payload["display_name"] == "성명 (1)"
        assert payload["fields"] == schema["fields"]

        duplicated_again = client.post(f"/api/schemas/{schema['id']}/duplicate")
        assert duplicated_again.status_code == 200, duplicated_again.text
        assert duplicated_again.json()["name"] == "성명 (2)"


def test_schema_recommendation_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            document = upload_png(client)
            response = client.post("/api/schemas/recommendations", json={"document_id": document["document_id"]})
            assert response.status_code == 200, response.text
            payload = response.json()
            assert payload["name"] == "ai_recommended_schema"
            assert len(payload["fields"]) >= 3
            assert {field["output_format"] for field in payload["fields"]} <= {"string", "float", "date", "bool"}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_required_field_checklist_recommendation_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            document = upload_png(client)
            response = client.post(
                "/api/required-field-checklists/recommendations",
                json={"document_id": document["document_id"]},
            )
            assert response.status_code == 200, response.text
            payload = response.json()
            assert payload["name"] == "ai_recommended_checklist"
            assert len(payload["items"]) >= 3
            assert {item["evidence_type"] for item in payload["items"]} <= {
                "text_or_handwriting",
                "checkbox",
                "signature_or_stamp",
                "visual_mark",
                "other",
            }
            assert {item["region_id"] for item in payload["items"] if item["region_id"]} <= {
                region["id"] for region in payload["regions"]
            }
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_ai_draft_mock_mode_does_not_persist_sample_images() -> None:
    from app.database import SessionLocal, init_db

    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        init_db()
        db = SessionLocal()
        try:
            before_documents = db.query(Document).count()
            before_batches = db.query(Batch).count()
        finally:
            db.close()

        with get_client() as client:
            response = client.post(
                "/api/workflows/ai-draft",
                files=[
                    ("files", ("sample-one.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("sample-two.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            payload = response.json()
            assert payload["sample_count"] == 2
            assert payload["images_persisted"] is False
            assert payload["schema_draft"]["name"] == "ai_recommended_schema"
            assert len(payload["schema_draft"]["fields"]) >= 3
            assert payload["checklist_draft"]["name"] == "ai_recommended_checklist"
            assert len(payload["definition"]["nodes"]) == 5
            assert any(node["id"] == "ai_kie" for node in payload["definition"]["nodes"])
            assert any(node["id"] == "ai_required" for node in payload["definition"]["nodes"])

        db = SessionLocal()
        try:
            assert db.query(Document).count() == before_documents
            assert db.query(Batch).count() == before_batches
        finally:
            db.close()
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_ai_draft_rejects_more_than_ten_images() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            response = client.post(
                "/api/workflows/ai-draft",
                files=[
                    ("files", (f"sample-{index}.png", ONE_BY_ONE_PNG, "image/png"))
                    for index in range(11)
                ],
            )
            assert response.status_code == 413, response.text
            assert "up to 10" in response.json()["detail"]
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_ai_draft_can_be_edited_and_saved_as_workflow() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        unique = uuid4().hex[:8]
        with get_client() as client:
            draft_response = client.post(
                "/api/workflows/ai-draft",
                files=[("files", ("sample.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert draft_response.status_code == 200, draft_response.text
            draft = draft_response.json()

            schema_payload = draft["schema_draft"]
            schema_payload["name"] = f"edited_ai_schema_{unique}"
            schema_payload["display_name"] = "Edited AI Schema"
            schema_payload["fields"][0]["key_name"] = "edited_customer_name"
            schema_payload["fields"][0]["description"] = "Customer name edited inside Workflow Builder."
            schema_payload["fields"][0]["judgement_enabled"] = True
            schema_response = client.post("/api/schemas", json=schema_payload)
            assert schema_response.status_code == 200, schema_response.text
            schema = schema_response.json()
            assert schema["fields"][0]["key_name"] == "edited_customer_name"
            assert schema["fields"][0]["judgement_enabled"] is True

            checklist_payload = draft["checklist_draft"]
            assert checklist_payload is not None
            checklist_payload["name"] = f"edited_ai_checklist_{unique}"
            checklist_payload["items"][0]["item_name"] = "edited_required_item"
            checklist_response = client.post("/api/required-field-checklists", json=checklist_payload)
            assert checklist_response.status_code == 200, checklist_response.text
            checklist = checklist_response.json()
            assert checklist["items"][0]["item_name"] == "edited_required_item"

            definition = draft["definition"]
            for node in definition["nodes"]:
                if node["data"]["kind"] == "kie":
                    node["data"]["config"] = {"schema_id": schema["id"]}
                if node["data"]["kind"] == "required-checker":
                    node["data"]["config"] = {"checklist_id": checklist["id"]}

            workflow_response = client.post(
                "/api/workflows",
                json={"name": f"edited_ai_workflow_{unique}", "definition": definition},
            )
            assert workflow_response.status_code == 200, workflow_response.text
            workflow = workflow_response.json()
            assert workflow["validation_warnings"] == []
            kie_node = next(node for node in workflow["definition"]["nodes"] if node["data"]["kind"] == "kie")
            required_node = next(node for node in workflow["definition"]["nodes"] if node["data"]["kind"] == "required-checker")
            assert kie_node["data"]["config"]["schema_id"] == schema["id"]
            assert required_node["data"]["config"]["checklist_id"] == checklist["id"]
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_schema_description_recommendation_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            response = client.post(
                "/api/schemas/description-recommendations",
                json={
                    "name": "consent_schema",
                    "current_description": "Old description",
                    "fields": [
                        {
                            "key_name": "본인 성명",
                            "description": "문서 하단 서명 영역의 본인 성명",
                            "output_format": "string",
                        },
                        {
                            "key_name": "동의 여부",
                            "description": "체크박스 선택 상태를 기준으로 한 동의 여부",
                            "output_format": "bool",
                        },
                    ],
                },
            )
            assert response.status_code == 200, response.text
            payload = response.json()
            assert "consent_schema" in payload["description"]
            assert "본인 성명" in payload["description"]
            assert payload["reasoning"]
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_extraction_mock_mode_returns_evidence_and_normalized_values() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            document = upload_png(client)
            schema = create_schema(client)
            job_response = client.post(
                "/api/extraction-jobs",
                json={"document_id": document["document_id"], "schema_id": schema["id"]},
            )
            assert job_response.status_code == 200, job_response.text
            job_id = job_response.json()["job_id"]

            job = client.get(f"/api/extraction-jobs/{job_id}").json()
            assert job["status"] == "completed"
            values = job["result"]["validated_output"]["values"]
            assert values["invoice_number"]["page"] == 1
            assert values["invoice_number"]["evidence"]
            assert values["total_amount"]["normalized_value"] == 1234.5

            jobs = client.get(f"/api/extraction-jobs?document_id={document['document_id']}").json()
            assert any(item["job_id"] == job_id for item in jobs)

            csv_export = client.get(f"/api/extraction-results/{job['result_id']}/export?format=csv")
            assert csv_export.status_code == 200
            assert csv_export.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_export.headers["content-type"]
            assert "evidence" in csv_export.text.splitlines()[0]

            xlsx_export = client.get(f"/api/extraction-results/{job['result_id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_export)

            corrected_output = job["result"]["validated_output"]
            corrected_output["values"]["invoice_number"]["value"] = "INV-EDITED"
            patch = client.patch(
                f"/api/extraction-results/{job['result_id']}",
                json={"corrected_output": corrected_output},
            )
            assert patch.status_code == 200, patch.text
            assert patch.json()["corrected_output"]["values"]["invoice_number"]["value"] == "INV-EDITED"
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_draft_extraction_does_not_list_schema() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            document = upload_png(client)
            response = client.post(
                "/api/extraction-jobs/draft",
                json={
                    "document_id": document["document_id"],
                    "schema": {
                        "name": "unsaved_draft_schema",
                        "display_name": "Unsaved Draft Schema",
                        "fields": [
                            {
                                "key_name": "draft_value",
                                "description": "Value visible in the draft document.",
                                "output_format": "string",
                            }
                        ],
                    },
                },
            )
            assert response.status_code == 200, response.text
            job = client.get(f"/api/extraction-jobs/{response.json()['job_id']}").json()
            assert job["status"] == "completed"
            assert job["result"]["validated_output"]["values"]["draft_value"]["value"] == "Sample draft_value"

            schemas = client.get("/api/schemas").json()
            assert all(item["name"] != "unsaved_draft_schema" for item in schemas)
            hidden_schema = client.get(f"/api/schemas/{job['schema_id']}")
            assert hidden_schema.status_code == 200
            assert hidden_schema.json()["ephemeral"] is True
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_schema_name_conflict_and_clear_parsing_history() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="conflict_schema")
            duplicate = client.post(
                "/api/schemas",
                json={
                    "name": "conflict_schema",
                    "fields": [
                        {
                            "key_name": "other",
                            "description": "Other field.",
                            "output_format": "string",
                        }
                    ],
                },
            )
            assert duplicate.status_code == 409

            document = upload_png(client)
            job_response = client.post(
                "/api/extraction-jobs",
                json={"document_id": document["document_id"], "schema_id": schema["id"]},
            )
            assert job_response.status_code == 200
            job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()
            assert job["status"] == "completed"

            cleared = client.delete("/api/maintenance/parsing-history")
            assert cleared.status_code == 200, cleared.text
            payload = cleared.json()
            assert payload["status"] == "cleared"
            assert payload["counts"]["documents"] >= 1
            assert client.get("/api/documents").json() == []
            assert client.get("/api/extraction-jobs").json() == []
            assert any(item["id"] == schema["id"] for item in client.get("/api/schemas").json())
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_extraction_low_confidence_requires_review(monkeypatch) -> None:
    def fake_extract(fields, image_paths=None, image_inputs=None):
        return {
            "invoice_number": {
                "value": "INV-LOW",
                "page": 1,
                "evidence": "low confidence evidence",
                "confidence": 0.0,
            },
            "total_amount": {
                "value": "10.00",
                "page": 1,
                "evidence": "high confidence evidence",
                "confidence": 0.95,
            },
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)

    with get_client() as client:
        document = upload_png(client)
        schema = create_schema(client)
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "needs_review"
    result = job["result"]
    assert result["validated_output"]["status"] == "needs_review"
    values = result["validated_output"]["values"]
    assert "low_confidence" in values["invoice_number"]["warnings"]
    assert values["total_amount"]["warnings"] == []
    assert "invoice_number:low_confidence" in result["validation_warnings"]


def test_extraction_null_without_evidence_is_not_detected(monkeypatch) -> None:
    def fake_extract(fields, image_paths=None, image_inputs=None):
        return {
            "invoice_number": {
                "value": None,
                "page": 1,
                "evidence": None,
                "confidence": 0.0,
            },
            "total_amount": {
                "value": "10.00",
                "page": 1,
                "evidence": "high confidence evidence",
                "confidence": 0.95,
            },
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)

    with get_client() as client:
        document = upload_png(client)
        schema = create_schema(client)
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "needs_review"
    values = job["result"]["validated_output"]["values"]
    assert values["invoice_number"]["value"] is None
    assert values["invoice_number"]["confidence"] is None
    assert "not_detected" in values["invoice_number"]["warnings"]
    assert "low_confidence" not in values["invoice_number"]["warnings"]
    assert "invoice_number:not_detected" in job["result"]["validation_warnings"]


def test_extraction_releases_db_connection_during_vlm_call(monkeypatch) -> None:
    from app.database import engine

    checked_out_counts: list[int] = []

    def fake_extract(fields, image_paths=None, image_inputs=None):
        checkedout = getattr(engine.pool, "checkedout", None)
        checked_out_counts.append(checkedout() if checkedout else 0)
        return {
            "invoice_number": {"value": "INV-POOL-001", "page": 1, "evidence": "test", "confidence": 0.9},
            "total_amount": {"value": "10.00", "page": 1, "evidence": "test", "confidence": 0.9},
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)

    with get_client() as client:
        document = upload_png(client)
        schema = create_schema(client)
        response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema["id"]},
        )
        assert response.status_code == 200, response.text
        job_id = response.json()["job_id"]
        job = client.get(f"/api/extraction-jobs/{job_id}").json()
        assert job["status"] == "completed"

    assert checked_out_counts
    assert max(checked_out_counts) == 0


def test_extraction_uses_schema_regions_for_cropped_inputs(monkeypatch) -> None:
    captured_calls: list[dict[str, object]] = []

    def fake_extract(fields, image_paths=None, image_inputs=None):
        captured_calls.append({"fields": fields, "image_paths": image_paths, "image_inputs": image_inputs})
        field_names = {field.key_name for field in fields}
        values = {}
        if "handwritten_name" in field_names:
            values["handwritten_name"] = {"value": "홍길동", "page": 1, "evidence": "region crop", "confidence": 0.91}
        if "handwritten_phone" in field_names:
            values["handwritten_phone"] = {"value": "010-0000-0000", "page": 1, "evidence": "region crop", "confidence": 0.9}
        if "document_date" in field_names:
            values["document_date"] = {"value": "2026.05.20", "page": 1, "evidence": "full page", "confidence": 0.9}
        return {
            **values,
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)

    with get_client() as client:
        document = client.post(
            "/api/documents",
            files={"file": ("sample.pdf", make_pdf_bytes(), "application/pdf")},
        ).json()
        schema = client.post(
            "/api/schemas",
            json={
                "name": "mixed_region_schema",
                "regions": [
                    {"id": "region_1", "name": "Handwriting block", "page": 1, "x": 0.1, "y": 0.1, "width": 0.5, "height": 0.4}
                ],
                "fields": [
                    {
                        "key_name": "handwritten_name",
                        "description": "손글씨 이름 영역",
                        "output_format": "string",
                        "region_id": "region_1",
                    },
                    {
                        "key_name": "handwritten_phone",
                        "description": "손글씨 연락처 영역",
                        "output_format": "string",
                        "region_id": "region_1",
                    },
                    {
                        "key_name": "document_date",
                        "description": "문서 전체에서 날짜",
                        "output_format": "date",
                    },
                ],
            },
        ).json()
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "completed"
    assert len(captured_calls) == 2

    full_page_call = next(call for call in captured_calls if [field.key_name for field in call["fields"]] == ["document_date"])
    full_page_inputs = full_page_call["image_inputs"]
    assert isinstance(full_page_inputs, list)
    assert any("Full document page 1" in item["label"] for item in full_page_inputs)

    region_call = next(call for call in captured_calls if {field.key_name for field in call["fields"]} == {"handwritten_name", "handwritten_phone"})
    region_inputs = region_call["image_inputs"]
    assert isinstance(region_inputs, list)
    assert any("Full page context" in item["label"] and "Handwriting block" in item["label"] for item in region_inputs)
    assert any("Masked full page context" in item["label"] and "Handwriting block" in item["label"] for item in region_inputs)
    assert any("Cropped extraction region" in item["label"] and "Handwriting block" in item["label"] for item in region_inputs)
    assert any("handwritten_name, handwritten_phone" in item["label"] for item in region_inputs)
    assert len(region_inputs) == 3
    region_fields = region_call["fields"]
    assert region_fields[0].region_id == "region_1"
    assert region_fields[1].region_id == "region_1"


def test_kie_prompts_are_split_by_region_presence() -> None:
    from app.prompts.kie import build_extraction_prompt, build_region_judgement_prompt
    from app.schemas import FieldDefinition

    full_prompt = build_extraction_prompt([
        FieldDefinition(key_name="document_date", description="Date visible on the document.", output_format="date")
    ])
    region_field = FieldDefinition(key_name="signature", description="Signature in the lower right area.", output_format="string", region_id="region_1")
    region_prompt = build_extraction_prompt([region_field])
    judgement_prompt = build_region_judgement_prompt(region_field, "서명", "first-stage evidence")

    assert "full document page images" in full_prompt
    assert "user-designated extraction region" not in full_prompt
    assert "labeled extraction region images" in region_prompt
    assert "crop image is already the user-designated extraction region" in region_prompt
    assert "masked image to confirm the region's original position" in region_prompt
    assert "verification step, not a re-extraction step" in judgement_prompt
    assert "default decision is judgement_status=correct" in judgement_prompt
    assert "If text such as 성명, 서명, 법정" in judgement_prompt


def test_kie_ai_judgement_corrects_enabled_field(monkeypatch) -> None:
    def fake_extract(fields, image_paths=None, image_inputs=None):
        return {
            field.key_name: {
                "value": "WRONG" if field.key_name == "invoice_number" else "1,234.50",
                "page": 1,
                "evidence": f"first evidence {field.key_name}",
                "confidence": 0.82,
            }
            for field in fields
        }

    judgement_calls = []
    correction_calls = []

    def fake_judge(field, initial_value, initial_evidence, image_inputs):
        judgement_calls.append((field.key_name, initial_value, initial_evidence, image_inputs))
        return {
            "judgement_status": "needs_correction",
            "reason": "The visible invoice number is different from the first-stage value.",
            "confidence": 0.93,
            "evidence": "Invoice number area",
        }

    def fake_correct(field, initial_value, initial_evidence, judgement_reason, image_inputs):
        correction_calls.append((field.key_name, judgement_reason, image_inputs))
        return {
            "value": "INV-2026-001",
            "page": 1,
            "evidence": "Corrected invoice number",
            "confidence": 0.95,
            "correction_reason": "Correct value is visible near the top.",
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)
    monkeypatch.setattr("app.extraction.judge_extraction_with_vlm", fake_judge)
    monkeypatch.setattr("app.extraction.correct_extraction_with_vlm", fake_correct)

    with get_client() as client:
        document = upload_png(client)
        schema = client.post(
            "/api/schemas",
            json={
                "name": "judgement_schema",
                "fields": [
                    {
                        "key_name": "invoice_number",
                        "description": "Invoice number near the top.",
                        "output_format": "string",
                        "judgement_enabled": True,
                    },
                    {
                        "key_name": "total_amount",
                        "description": "Final total amount.",
                        "output_format": "float",
                    },
                ],
            },
        )
        assert schema.status_code == 200, schema.text
        duplicated = client.post(f"/api/schemas/{schema.json()['id']}/duplicate")
        assert duplicated.status_code == 200, duplicated.text
        assert duplicated.json()["fields"][0]["judgement_enabled"] is True

        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema.json()["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "completed"
    values = job["result"]["validated_output"]["values"]
    assert values["invoice_number"]["value"] == "INV-2026-001"
    assert values["invoice_number"]["ai_review"]["corrected"] is True
    assert values["invoice_number"]["ai_review"]["initial_value"] == "WRONG"
    assert "ai_review" not in values["total_amount"]
    assert job["result"]["raw_model_output"]["invoice_number"]["value"] == "WRONG"
    assert [call[0] for call in judgement_calls] == ["invoice_number"]
    assert [call[0] for call in correction_calls] == ["invoice_number"]


def test_kie_ai_judgement_failure_needs_review_without_failing_job(monkeypatch) -> None:
    def fake_extract(fields, image_paths=None, image_inputs=None):
        return {
            field.key_name: {
                "value": "INV-1",
                "page": 1,
                "evidence": "first evidence",
                "confidence": 0.82,
            }
            for field in fields
        }

    def fake_judge(field, initial_value, initial_evidence, image_inputs):
        raise RuntimeError("synthetic judgement failure")

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)
    monkeypatch.setattr("app.extraction.judge_extraction_with_vlm", fake_judge)

    with get_client() as client:
        document = upload_png(client)
        schema = client.post(
            "/api/schemas",
            json={
                "name": "judgement_failure_schema",
                "fields": [
                    {
                        "key_name": "invoice_number",
                        "description": "Invoice number near the top.",
                        "output_format": "string",
                        "judgement_enabled": True,
                    }
                ],
            },
        )
        assert schema.status_code == 200, schema.text
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema.json()["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "needs_review"
    value = job["result"]["validated_output"]["values"]["invoice_number"]
    assert value["value"] == "INV-1"
    assert "ai_review_failed" in value["warnings"]
    assert value["ai_review"]["judgement_status"] == "failed"


@pytest.mark.parametrize(
    ("initial_value", "corrected_value", "expected_warning"),
    [
        ("성명", None, "ai_correction_discarded_null"),
        ("문어", "문이", "ai_correction_large_change"),
    ],
)
def test_kie_ai_correction_preserves_risky_korean_values(
    monkeypatch,
    initial_value: str,
    corrected_value: str | None,
    expected_warning: str,
) -> None:
    def fake_extract(fields, image_paths=None, image_inputs=None):
        return {
            field.key_name: {
                "value": initial_value,
                "page": 1,
                "evidence": f"first evidence {initial_value}",
                "confidence": 0.91,
            }
            for field in fields
        }

    def fake_judge(field, initial_value, initial_evidence, image_inputs):
        return {
            "judgement_status": "needs_correction",
            "reason": "Synthetic judgement requested a risky correction.",
            "confidence": 0.98,
            "evidence": "Synthetic evidence",
        }

    def fake_correct(field, initial_value, initial_evidence, judgement_reason, image_inputs):
        return {
            "value": corrected_value,
            "page": 1,
            "evidence": "Synthetic correction evidence",
            "confidence": 0.99,
            "correction_reason": "Synthetic risky correction.",
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)
    monkeypatch.setattr("app.extraction.judge_extraction_with_vlm", fake_judge)
    monkeypatch.setattr("app.extraction.correct_extraction_with_vlm", fake_correct)

    with get_client() as client:
        document = upload_png(client)
        schema = client.post(
            "/api/schemas",
            json={
                "name": f"judgement_risky_{expected_warning}",
                "fields": [
                    {
                        "key_name": "법정대리인성명",
                        "description": "Legal representative name in the handwritten field.",
                        "output_format": "string",
                        "judgement_enabled": True,
                    }
                ],
            },
        )
        assert schema.status_code == 200, schema.text
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema.json()["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "needs_review"
    value = job["result"]["validated_output"]["values"]["법정대리인성명"]
    assert value["value"] == initial_value
    assert expected_warning in value["warnings"]
    assert value["ai_review"]["judgement_status"] == "needs_correction"
    assert value["ai_review"]["corrected"] is False


def test_batch_cancel_marks_queued_jobs_canceled(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_batch_jobs", lambda batch_id, job_ids: None)

    with get_client() as client:
        schema = create_schema(client)
        response = client.post(
            "/api/batches",
            data={"schema_id": schema["id"]},
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert response.status_code == 200, response.text
        batch = response.json()
        assert batch["status"] == "running"
        assert batch["progress"] == 0

        canceled = client.post(f"/api/batches/{batch['id']}/cancel")
        assert canceled.status_code == 200, canceled.text
        payload = canceled.json()
        assert payload["status"] == "canceled"
        assert payload["canceled_count"] == 2
        assert payload["progress"] == 1
        assert payload["completed_at"] is not None
        assert {item["status"] for item in payload["items"]} == {"canceled"}


def test_batch_init_items_start_flow(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_batch_jobs", lambda batch_id, job_ids: None)

    with get_client() as client:
        schema = create_schema(client)
        initialized = client.post("/api/batches/init", json={"schema_id": schema["id"], "total_count": 2})
        assert initialized.status_code == 200, initialized.text
        batch = initialized.json()
        assert batch["status"] == "uploading"
        assert batch["uploaded_count"] == 0

        early_start = client.post(f"/api/batches/{batch['id']}/start")
        assert early_start.status_code == 422
        assert early_start.json()["detail"]["uploaded_count"] == 0

        partial = client.post(
            f"/api/batches/{batch['id']}/items",
            data={"client_file_ids": ["0:first.png:1:1"]},
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert partial.status_code == 200, partial.text
        partial_resume = client.post(f"/api/batches/{batch['id']}/resume")
        assert partial_resume.status_code == 422, partial_resume.text
        assert partial_resume.json()["detail"]["uploaded_count"] == 1
        assert partial_resume.json()["detail"]["total_count"] == 2

        appended = client.post(
            f"/api/batches/{batch['id']}/items",
            data={"client_file_ids": ["1:second.png:1:1"]},
            files=[("files", ("second.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert appended.status_code == 200, appended.text
        payload = appended.json()
        assert payload["uploaded_count"] == 2
        assert payload["queued_count"] == 2
        assert payload["preprocessing_count"] == 0

        summary = client.get(f"/api/batches/{batch['id']}/summary")
        assert summary.status_code == 200
        assert summary.json()["items"] == []
        assert summary.json()["queued_count"] == 2

        started = client.post(f"/api/batches/{batch['id']}/start")
        assert started.status_code == 200, started.text
        assert started.json()["status"] == "running"


def test_batch_export_csv_and_json_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client)
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            assert [item["filename"] for item in batch["items"]] == ["a_first.png", "z_last.png"]

            csv_response = client.get(f"/api/batches/{batch['id']}/export?format=csv")
            assert csv_response.status_code == 200, csv_response.text
            assert csv_response.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_response.headers["content-type"]
            csv_text = csv_response.text
            header = csv_text.splitlines()[0]
            assert "filename,document_id,job_id,status,error_message,invoice_number" in header
            assert "invoice_number_original" in header
            assert "invoice_number_ai_review_status" in header
            assert "total_amount_ai_corrected" in header
            assert "a_first.png" in csv_text
            assert "Sample invoice_number" in csv_text
            assert csv_text.index("a_first.png") < csv_text.index("z_last.png")

            json_response = client.get(f"/api/batches/{batch['id']}/export?format=json")
            assert json_response.status_code == 200, json_response.text
            payload = json_response.json()
            assert payload["batch_id"] == batch["id"]
            assert len(payload["rows"]) == 2
            assert [row["filename"] for row in payload["rows"]] == ["a_first.png", "z_last.png"]
            assert payload["rows"][0]["invoice_number"] == "Sample invoice_number"

            xlsx_response = client.get(f"/api/batches/{batch['id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_response)
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_export_job_batch_csv_download_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="async_export_batch_schema")
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[("files", ("invoice.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert response.status_code == 200, response.text
            batch = response.json()

            create_job = client.post(
                "/api/export-jobs",
                json={"owner_type": "batch", "owner_id": batch["id"], "format": "csv"},
            )
            assert create_job.status_code == 202, create_job.text
            queued = create_job.json()
            assert queued["owner_type"] == "batch"
            assert queued["status"] in {"queued", "running", "completed"}

            completed = wait_for_export_job(client, queued["id"])
            assert completed["status"] == "completed", completed
            assert completed["filename"].endswith(".csv")
            assert completed["size_bytes"] > 0

            download = client.get(f"/api/export-jobs/{queued['id']}/download")
            assert download.status_code == 200, download.text
            assert download.content.startswith(b"\xef\xbb\xbf")
            assert "invoice_number" in download.text.splitlines()[0]
            assert "invoice.png" in download.text

            listed = client.get(f"/api/export-jobs?owner_type=batch&owner_id={batch['id']}&limit=5")
            assert listed.status_code == 200, listed.text
            assert queued["id"] in {item["id"] for item in listed.json()}

            missing_owner = client.post(
                "/api/export-jobs",
                json={"owner_type": "batch", "owner_id": "batch_missing", "format": "csv"},
            )
            assert missing_owner.status_code == 404
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_export_job_retry_failed_job_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="retry_export_batch_schema")
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[("files", ("retry.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert response.status_code == 200, response.text
            batch = response.json()

            from app.database import SessionLocal

            db = SessionLocal()
            source = ExportJob(
                owner_type="batch",
                owner_id=batch["id"],
                format="csv",
                status="failed",
                error_message="network timeout",
                started_at=datetime.utcnow(),
                completed_at=datetime.utcnow(),
            )
            db.add(source)
            db.commit()
            source_job_id = source.id
            db.close()

            retry = client.post(f"/api/export-jobs/{source_job_id}/retry")
            assert retry.status_code == 202, retry.text
            retried = retry.json()
            assert retried["id"] != source_job_id
            assert retried["owner_id"] == batch["id"]

            completed = wait_for_export_job(client, retried["id"])
            assert completed["status"] == "completed", completed

            failed_jobs = client.get(f"/api/export-jobs?owner_type=batch&owner_id={batch['id']}&status=failed")
            assert failed_jobs.status_code == 200, failed_jobs.text
            assert [item["id"] for item in failed_jobs.json()] == [source_job_id]

            retry_completed = client.post(f"/api/export-jobs/{completed['id']}/retry")
            assert retry_completed.status_code == 409
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_export_job_worker_only_claims_queued_jobs() -> None:
    from app.database import SessionLocal, init_db
    from app.main import _run_export_job

    init_db()
    db = SessionLocal()
    job = ExportJob(
        owner_type="batch",
        owner_id="batch_missing",
        format="csv",
        status="running",
        started_at=datetime.utcnow(),
    )
    db.add(job)
    db.commit()
    job_id = job.id
    db.close()

    _run_export_job(job_id)

    db = SessionLocal()
    try:
        loaded = db.get(ExportJob, job_id)
        assert loaded is not None
        assert loaded.status == "running"
        assert loaded.error_message is None
    finally:
        db.close()


def test_export_worker_recovers_interrupted_running_job_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="recover_export_batch_schema")
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[("files", ("recover.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert response.status_code == 200, response.text
            batch = response.json()

        from app.database import SessionLocal
        from app.main import _reset_interrupted_export_jobs, _run_export_job

        db = SessionLocal()
        job = ExportJob(
            owner_type="batch",
            owner_id=batch["id"],
            format="csv",
            status="running",
            started_at=datetime.utcnow(),
        )
        db.add(job)
        db.commit()
        job_id = job.id
        db.close()

        _reset_interrupted_export_jobs()

        db = SessionLocal()
        try:
            recovered = db.get(ExportJob, job_id)
            assert recovered is not None
            assert recovered.status == "queued"
            assert recovered.error_message == "Recovered after server restart"
        finally:
            db.close()

        _run_export_job(job_id)

        db = SessionLocal()
        try:
            completed = db.get(ExportJob, job_id)
            assert completed is not None
            assert completed.status == "completed"
            assert completed.storage_path
            assert completed.size_bytes > 0
        finally:
            db.close()
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_batch_finalizes_after_mock_jobs() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client)
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()

            loaded = client.get(f"/api/batches/{batch['id']}")
            assert loaded.status_code == 200, loaded.text
            payload = loaded.json()
            assert payload["status"] == "completed"
            assert payload["progress"] == 1
            assert payload["completed_count"] == 2
            assert payload["completed_at"] is not None
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_batch_high_concurrent_request_count_does_not_exhaust_db_pool() -> None:
    previous_workers = os.environ.get("VLM_MAX_CONCURRENT_REQUESTS")
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = "16"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client)
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", (f"batch_{index}.png", ONE_BY_ONE_PNG, "image/png"))
                    for index in range(20)
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()

            loaded = client.get(f"/api/batches/{batch['id']}")
            assert loaded.status_code == 200, loaded.text
            payload = loaded.json()
            assert payload["status"] == "completed"
            assert payload["completed_count"] == 20
            assert payload["failed_count"] == 0

            recent = client.get("/api/batches?limit=12")
            assert recent.status_code == 200, recent.text
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        if previous_workers is None:
            os.environ.pop("VLM_MAX_CONCURRENT_REQUESTS", None)
        else:
            os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = previous_workers
        get_settings.cache_clear()


def test_document_classifier_config_single_job_and_patch() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            classifier = create_document_classifier(client)
            updated = client.patch(
                f"/api/document-classifiers/{classifier['id']}",
                json={
                    "description": "Updated classifier description",
                    "allow_unknown": True,
                },
            )
            assert updated.status_code == 200, updated.text
            assert updated.json()["description"] == "Updated classifier description"

            document = upload_png(client)
            response = client.post(
                "/api/classification-jobs",
                json={"document_id": document["document_id"], "classifier_id": classifier["id"]},
            )
            assert response.status_code == 200, response.text
            job = client.get(f"/api/classification-jobs/{response.json()['job_id']}").json()
            assert job["status"] == "completed"
            output = job["result"]["validated_output"]
            assert output["status"] == "classified"
            assert output["class_name"] == "contract"
            assert output["confidence"] == 0.88

            corrected = {**output, "status": "unknown", "class_name": None}
            patch = client.patch(
                f"/api/classification-results/{job['result_id']}",
                json={"corrected_output": corrected, "reviewed": True},
            )
            assert patch.status_code == 200, patch.text
            assert patch.json()["corrected_output"]["status"] == "unknown"

            deleted = client.delete(f"/api/document-classifiers/{classifier['id']}")
            assert deleted.status_code == 200
            assert deleted.json()["archived"] is True
            assert all(item["id"] != classifier["id"] for item in client.get("/api/document-classifiers").json())
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_document_classifier_duplicate_copies_classes() -> None:
    with get_client() as client:
        classifier = create_document_classifier(client, name="신청서 분류")
        duplicated = client.post(f"/api/document-classifiers/{classifier['id']}/duplicate")
        assert duplicated.status_code == 200, duplicated.text
        payload = duplicated.json()
        assert payload["id"] != classifier["id"]
        assert payload["name"] == "신청서 분류 (1)"
        assert payload["classes"] == classifier["classes"]
        assert payload["allow_unknown"] == classifier["allow_unknown"]


def test_document_classifier_batch_cancel_and_export(monkeypatch) -> None:
    with monkeypatch.context() as patch_context:
        patch_context.setattr("app.main.run_classification_batch", lambda batch_id, job_ids: None)

        with get_client() as client:
            classifier = create_document_classifier(client)
            response = client.post(
                "/api/classification-batches",
                data={"classifier_id": classifier["id"]},
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            assert [item["filename"] for item in batch["items"]] == ["a_first.png", "z_last.png"]

            canceled = client.post(f"/api/classification-batches/{batch['id']}/cancel")
            assert canceled.status_code == 200, canceled.text
            assert canceled.json()["status"] == "canceled"
            assert canceled.json()["canceled_count"] == 2
            assert canceled.json()["completed_at"] is not None

    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            classifier = create_document_classifier(client, name="batch_classifier")
            response = client.post(
                "/api/classification-batches",
                data={"classifier_id": classifier["id"]},
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            loaded = client.get(f"/api/classification-batches/{batch['id']}").json()
            assert loaded["status"] == "completed"
            assert loaded["completed_count"] == 2

            csv_response = client.get(f"/api/classification-batches/{batch['id']}/export?format=csv")
            assert csv_response.status_code == 200, csv_response.text
            assert csv_response.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_response.headers["content-type"]
            assert "classification_status,class_name,confidence,reason,evidence" in csv_response.text.splitlines()[0]
            assert csv_response.text.index("a_first.png") < csv_response.text.index("z_last.png")

            json_response = client.get(f"/api/classification-batches/{batch['id']}/export?format=json")
            assert json_response.status_code == 200
            rows = json_response.json()["rows"]
            assert [row["filename"] for row in rows] == ["a_first.png", "z_last.png"]
            assert rows[0]["class_name"] == "contract"

            xlsx_response = client.get(f"/api/classification-batches/{batch['id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_response)
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_required_field_checklist_single_job_and_region_validation() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            invalid = client.post(
                "/api/required-field-checklists",
                json={
                    "name": "invalid_checklist",
                    "regions": [],
                    "items": [
                        {
                            "item_name": "서명",
                            "description": "서명 존재 여부",
                            "evidence_type": "signature_or_stamp",
                            "required": True,
                            "region_id": "missing_region",
                        }
                    ],
                },
            )
            assert invalid.status_code == 422

            custom = client.post(
                "/api/required-field-checklists",
                json={
                    "name": "custom_evidence_checklist",
                    "regions": [],
                    "items": [
                        {
                            "item_name": "수기 메모",
                            "description": "사용자가 직접 정의한 증거 유형을 확인합니다.",
                            "evidence_type": "수기 메모/특이사항",
                            "required": True,
                        }
                    ],
                },
            )
            assert custom.status_code == 200, custom.text
            assert custom.json()["items"][0]["evidence_type"] == "수기 메모/특이사항"

            checklist = create_required_field_checklist(client)
            document = upload_png(client)
            response = client.post(
                "/api/required-field-check-jobs",
                json={"document_id": document["document_id"], "checklist_id": checklist["id"]},
            )
            assert response.status_code == 200, response.text
            job = client.get(f"/api/required-field-check-jobs/{response.json()['job_id']}").json()
            assert job["status"] == "needs_review"
            output = job["result"]["validated_output"]
            assert output["overall_status"] == "needs_review"
            assert [item["item_name"] for item in output["items"]] == ["성명", "서명", "체크박스"]
            assert output["items"][0]["status"] == "present"
            assert output["items"][2]["status"] == "uncertain"

            corrected = {
                **output,
                "overall_status": "complete",
                "items": [{**item, "status": "present"} for item in output["items"]],
            }
            patch = client.patch(
                f"/api/required-field-check-results/{job['result_id']}",
                json={"corrected_output": corrected, "reviewed": True},
            )
            assert patch.status_code == 200, patch.text
            assert patch.json()["corrected_output"]["overall_status"] == "complete"

            deleted = client.delete(f"/api/required-field-checklists/{checklist['id']}")
            assert deleted.status_code == 200
            assert deleted.json()["archived"] is True
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_required_checker_accepts_pydantic_item_object_output() -> None:
    from app.document_modules import RequiredFieldContext, _validate_required_field_output
    from app.extraction import DocumentSnapshot
    from app.schemas import RequiredFieldItem

    context = RequiredFieldContext(
        document=DocumentSnapshot(id="doc_1", storage_path="", pages=[]),
        checklist_id="checklist_1",
        items=[
            RequiredFieldItem(item_name="서명", description="서명 확인", required=True),
            RequiredFieldItem(item_name="작성일", description="작성일 확인", required=True),
        ],
        regions=[],
    )
    raw_values = {
        "overall_status": "needs_review",
        "items": {
            "서명": {"status": "present", "confidence": 0.9, "evidence": "서명이 보입니다.", "page": 1},
            "작성일": {"status": "missing", "confidence": 0.8, "evidence": None, "page": None},
        },
    }

    validated = _validate_required_field_output(raw_values, context)

    assert validated["overall_status"] == "incomplete"
    assert [item["item_name"] for item in validated["items"]] == ["서명", "작성일"]
    assert validated["items"][0]["status"] == "present"
    assert validated["items"][1]["status"] == "missing"


def test_required_checker_splits_full_page_and_region_requests(monkeypatch) -> None:
    captured_calls: list[dict[str, object]] = []

    def fake_check(items, regions, image_paths=None, image_inputs=None):
        captured_calls.append(
            {
                "item_names": [item.item_name for item in items],
                "region_ids": [region.id for region in regions],
                "labels": [item["label"] for item in image_inputs or []],
            }
        )
        return {
            "overall_status": "complete",
            "items": [
                {
                    "item_name": item.item_name,
                    "status": "present",
                    "confidence": 0.9,
                    "evidence": "captured",
                    "page": 1,
                }
                for item in items
            ],
        }

    monkeypatch.setattr("app.document_modules.check_required_fields_with_vlm", fake_check)

    with get_client() as client:
        checklist = create_required_field_checklist(client, name="split_required_checklist")
        document = upload_png(client)
        response = client.post(
            "/api/required-field-check-jobs",
            json={"document_id": document["document_id"], "checklist_id": checklist["id"]},
        )
        assert response.status_code == 200, response.text
        job = client.get(f"/api/required-field-check-jobs/{response.json()['job_id']}").json()

    assert job["status"] == "completed"
    assert len(captured_calls) == 2
    full_page_call = next(call for call in captured_calls if call["region_ids"] == [])
    assert full_page_call["item_names"] == ["성명", "체크박스"]
    assert any("Full document page" in label for label in full_page_call["labels"])

    region_call = next(call for call in captured_calls if call["region_ids"] == ["signature_region"])
    assert region_call["item_names"] == ["서명"]
    assert any("Masked context" in label and "서명 영역" in label for label in region_call["labels"])
    assert any("Cropped required field region" in label and "서명 영역" in label for label in region_call["labels"])
    assert not any("Full document page" in label for label in region_call["labels"])


def test_required_field_checklist_duplicate_copies_items_and_regions() -> None:
    with get_client() as client:
        checklist = create_required_field_checklist(client, name="필수 정보")
        duplicated = client.post(f"/api/required-field-checklists/{checklist['id']}/duplicate")
        assert duplicated.status_code == 200, duplicated.text
        payload = duplicated.json()
        assert payload["id"] != checklist["id"]
        assert payload["name"] == "필수 정보 (1)"
        assert payload["items"] == checklist["items"]
        assert payload["regions"] == checklist["regions"]


def test_required_field_check_batch_export_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            checklist = create_required_field_checklist(client, name="batch_checklist")
            response = client.post(
                "/api/required-field-check-batches",
                data={"checklist_id": checklist["id"]},
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            assert [item["filename"] for item in batch["items"]] == ["a_first.png", "z_last.png"]

            loaded = client.get(f"/api/required-field-check-batches/{batch['id']}").json()
            assert loaded["status"] == "completed"
            assert loaded["completed_count"] == 2

            csv_response = client.get(f"/api/required-field-check-batches/{batch['id']}/export?format=csv")
            assert csv_response.status_code == 200, csv_response.text
            assert csv_response.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_response.headers["content-type"]
            header = csv_response.text.splitlines()[0]
            assert "overall_status" in header
            assert "성명_status" in header
            assert csv_response.text.index("a_first.png") < csv_response.text.index("z_last.png")

            json_response = client.get(f"/api/required-field-check-batches/{batch['id']}/export?format=json")
            assert json_response.status_code == 200
            rows = json_response.json()["rows"]
            assert rows[0]["성명_status"] == "present"

            xlsx_response = client.get(f"/api/required-field-check-batches/{batch['id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_response)
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_required_field_check_batch_cancel_marks_queued_jobs_canceled(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_required_field_check_batch", lambda batch_id, job_ids: None)

    with get_client() as client:
        checklist = create_required_field_checklist(client, name="cancel_checklist")
        response = client.post(
            "/api/required-field-check-batches",
            data={"checklist_id": checklist["id"]},
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert response.status_code == 200, response.text
        batch = response.json()
        assert batch["status"] == "running"

        canceled = client.post(f"/api/required-field-check-batches/{batch['id']}/cancel")
        assert canceled.status_code == 200, canceled.text
        payload = canceled.json()
        assert payload["status"] == "canceled"
        assert payload["canceled_count"] == 2
        assert payload["progress"] == 1
        assert payload["completed_at"] is not None
        assert {item["status"] for item in payload["items"]} == {"canceled"}


def test_workflow_definition_validation_and_branch_run_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_invoice_schema")
            classifier = create_document_classifier(client, name="workflow_classifier")
            checklist = create_required_field_checklist(client, name="workflow_checklist")

            invalid = client.post(
                "/api/workflows",
                json={
                    "name": "invalid_workflow",
                    "definition": {
                        "nodes": [
                            {"id": "input", "data": {"kind": "input"}},
                            {"id": "classifier", "data": {"kind": "classifier", "config": {}}},
                            {"id": "export", "data": {"kind": "export"}},
                        ],
                        "edges": [
                            {"id": "e1", "source": "input", "target": "classifier"},
                            {"id": "e2", "source": "classifier", "target": "export"},
                        ],
                    },
                },
            )
            assert invalid.status_code == 422
            assert "Classifier node classifier" in invalid.text

            workflow = client.post(
                "/api/workflows",
                json={
                    "name": "분기 워크플로우",
                    "description": "Classifier 결과에 따라 계약서 경로를 실행합니다.",
                    "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"]),
                },
            )
            assert workflow.status_code == 200, workflow.text
            assert workflow.json()["validation_warnings"] == []

            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            assert run["status"] in {"completed", "needs_review"}
            assert run["total_count"] == 2
            assert [item["filename"] for item in run["items"]] == ["a_first.png", "z_last.png"]
            first = run["items"][0]["result"]
            assert first["classification"]["class_name"] == "contract"
            assert first["branch_path"] == "class:contract"
            assert "invoice_number" in first["kie_values"]
            assert first["required_items"]["성명"]["status"] == "present"

            csv_response = client.get(f"/api/workflow-runs/{run['id']}/export?format=csv")
            assert csv_response.status_code == 200, csv_response.text
            assert csv_response.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_response.headers["content-type"]
            header = csv_response.text.splitlines()[0]
            assert "classification_status" in header
            assert "upload_duration_ms" in header
            assert "inference_duration_ms" in header
            assert "kie_invoice_number" in header
            assert "required_성명_status" in header
            assert csv_response.text.index("a_first.png") < csv_response.text.index("z_last.png")

            json_response = client.get(f"/api/workflow-runs/{run['id']}/export?format=json")
            assert json_response.status_code == 200
            assert isinstance(run["upload_duration_ms"], int)
            assert isinstance(run["inference_duration_ms"], int)
            assert all(isinstance(item["upload_duration_ms"], int) for item in run["items"])
            assert all(isinstance(item["inference_duration_ms"], int) for item in run["items"])
            assert json_response.json()["rows"][0]["upload_duration_ms"] is not None
            assert json_response.json()["rows"][0]["class_name"] == "contract"

            xlsx_response = client.get(f"/api/workflow-runs/{run['id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_response)

            export_job = client.post(
                "/api/export-jobs",
                json={"owner_type": "workflow_run", "owner_id": run["id"], "format": "json"},
            )
            assert export_job.status_code == 202, export_job.text
            completed_export = wait_for_export_job(client, export_job.json()["id"])
            assert completed_export["status"] == "completed", completed_export
            async_download = client.get(f"/api/export-jobs/{completed_export['id']}/download")
            assert async_download.status_code == 200, async_download.text
            assert async_download.json()["workflow_run_id"] == run["id"]
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_summary_polling_uses_aggregate_counts_without_result_payloads() -> None:
    from sqlalchemy import event

    from app.database import SessionLocal, engine, init_db
    from app.models import WorkflowDefinition, WorkflowRunItem

    init_db()
    db = SessionLocal()
    try:
        workflow = WorkflowDefinition(
            name="summary_polling_workflow",
            definition_json=json.dumps({"nodes": [], "edges": []}),
        )
        db.add(workflow)
        db.flush()
        run = WorkflowRun(
            workflow_id=workflow.id,
            workflow_name=workflow.name,
            status="running",
            total_count=4,
        )
        db.add(run)
        db.flush()
        for index, status in enumerate(["completed", "needs_review", "failed", "canceled"]):
            db.add(
                WorkflowRunItem(
                    run_id=run.id,
                    document_id=f"doc_summary_{index}",
                    filename=f"summary_{index}.png",
                    upload_index=index,
                    status=status,
                    result_json=json.dumps({"payload": "x" * 10000}),
                )
            )
        db.commit()
        run_id = run.id
    finally:
        db.close()

    item_queries: list[str] = []

    def capture_item_queries(_conn, _cursor, statement, _parameters, _context, _executemany) -> None:
        normalized = " ".join(statement.lower().split())
        if "workflow_run_items" in normalized:
            item_queries.append(normalized)

    event.listen(engine, "before_cursor_execute", capture_item_queries)
    try:
        with get_client() as client:
            summary_response = client.get(f"/api/workflow-runs/{run_id}/summary")
            list_response = client.get("/api/workflow-runs?limit=100")
    finally:
        event.remove(engine, "before_cursor_execute", capture_item_queries)

    assert summary_response.status_code == 200, summary_response.text
    summary = summary_response.json()
    assert summary["items"] == []
    assert summary["status"] == "completed_with_errors"
    assert summary["completed_count"] == 1
    assert summary["failed_count"] == 1
    assert summary["needs_review_count"] == 1
    assert summary["canceled_count"] == 1
    assert summary["progress"] == 1

    assert list_response.status_code == 200, list_response.text
    listed = next(item for item in list_response.json() if item["id"] == run_id)
    assert listed["items"] == []
    assert listed["failed_count"] == 1

    assert item_queries
    assert any("count" in query and "group by" in query for query in item_queries)
    assert all("result_json" not in query for query in item_queries)


def test_workflow_summary_polling_handles_1000_item_run_with_stable_counts() -> None:
    from sqlalchemy import event

    from app.database import SessionLocal, engine, init_db
    from app.models import WorkflowDefinition, WorkflowRunItem

    init_db()
    status_sequence = (
        ["completed"] * 700
        + ["needs_review"] * 100
        + ["failed"] * 50
        + ["canceled"] * 25
        + ["running"] * 75
        + ["waiting_for_document"] * 50
    )
    db = SessionLocal()
    try:
        workflow = WorkflowDefinition(
            name="summary_polling_1000_workflow",
            definition_json=json.dumps({"nodes": [], "edges": []}),
        )
        db.add(workflow)
        db.flush()
        run = WorkflowRun(
            workflow_id=workflow.id,
            workflow_name=workflow.name,
            status="running",
            total_count=len(status_sequence),
        )
        db.add(run)
        db.flush()
        db.add_all(
            [
                WorkflowRunItem(
                    run_id=run.id,
                    document_id=f"doc_summary_1000_{index}",
                    filename=f"summary_1000_{index:04d}.png",
                    upload_index=index,
                    status=status,
                    result_json=json.dumps({"payload": "x" * 100}),
                )
                for index, status in enumerate(status_sequence)
            ]
        )
        db.commit()
        run_id = run.id
    finally:
        db.close()

    item_queries: list[str] = []

    def capture_item_queries(_conn, _cursor, statement, _parameters, _context, _executemany) -> None:
        normalized = " ".join(statement.lower().split())
        if "workflow_run_items" in normalized:
            item_queries.append(normalized)

    event.listen(engine, "before_cursor_execute", capture_item_queries)
    try:
        with get_client() as client:
            summary_response = client.get(f"/api/workflow-runs/{run_id}/summary")
    finally:
        event.remove(engine, "before_cursor_execute", capture_item_queries)

    assert summary_response.status_code == 200, summary_response.text
    summary = summary_response.json()
    assert summary["items"] == []
    assert summary["total_count"] == 1000
    assert summary["uploaded_count"] == 1000
    assert summary["completed_count"] == 700
    assert summary["needs_review_count"] == 100
    assert summary["failed_count"] == 50
    assert summary["canceled_count"] == 25
    assert summary["running_count"] == 75
    assert summary["preprocessing_count"] == 50
    assert summary["progress_phase"] == "preprocessing"
    assert summary["progress"] == 0.875
    assert {"vlm_active_count", "vlm_waiting_count", "vlm_limit"}.issubset(summary)
    assert item_queries
    assert any("count" in query and "group by" in query for query in item_queries)
    assert all("result_json" not in query for query in item_queries)


@pytest.mark.parametrize(
    "case",
    [
        {
            "owner_model": Batch,
            "item_model_name": "BatchItem",
            "job_model_name": "ExtractionJob",
            "owner_kwargs": {"schema_id": "schema_summary_polling", "schema_version": 1},
            "item_owner_field": "batch_id",
            "job_kwargs": {"schema_id": "schema_summary_polling", "schema_version": 1},
            "summary_path": "/api/batches/{id}/summary",
            "list_path": "/api/batches?limit=100&include_items=false",
            "item_table": "batch_items",
        },
        {
            "owner_model": ClassificationBatch,
            "item_model_name": "ClassificationBatchItem",
            "job_model_name": "ClassificationJob",
            "owner_kwargs": {"classifier_id": "classifier_summary_polling"},
            "item_owner_field": "batch_id",
            "job_kwargs": {"classifier_id": "classifier_summary_polling"},
            "summary_path": "/api/classification-batches/{id}/summary",
            "list_path": "/api/classification-batches?limit=100&include_items=false",
            "item_table": "classification_batch_items",
        },
        {
            "owner_model": RequiredFieldCheckBatch,
            "item_model_name": "RequiredFieldCheckBatchItem",
            "job_model_name": "RequiredFieldCheckJob",
            "owner_kwargs": {"checklist_id": "checklist_summary_polling"},
            "item_owner_field": "batch_id",
            "job_kwargs": {"checklist_id": "checklist_summary_polling"},
            "summary_path": "/api/required-field-check-batches/{id}/summary",
            "list_path": "/api/required-field-check-batches?limit=100&include_items=false",
            "item_table": "required_field_check_batch_items",
        },
    ],
)
def test_module_batch_summary_and_light_list_use_aggregate_counts(case: dict[str, Any]) -> None:
    from sqlalchemy import event

    from app import models as model_module
    from app.database import SessionLocal, engine, init_db

    init_db()
    item_model = getattr(model_module, case["item_model_name"])
    job_model = getattr(model_module, case["job_model_name"])
    statuses = ["completed", "needs_review", "failed", "canceled", "running", "waiting_for_document"]

    db = SessionLocal()
    try:
        owner = case["owner_model"](status="running", total_count=len(statuses), **case["owner_kwargs"])
        db.add(owner)
        db.flush()
        for index, status in enumerate(statuses):
            job = job_model(
                document_id=f"doc_module_summary_{case['item_table']}_{index}",
                status=status,
                error_message="failed by test" if status == "failed" else None,
                **case["job_kwargs"],
            )
            db.add(job)
            db.flush()
            db.add(
                item_model(
                    **{case["item_owner_field"]: owner.id},
                    document_id=job.document_id,
                    job_id=job.id,
                    filename=f"module_summary_{index}.png",
                    upload_index=index,
                    client_file_id=f"test:{index}",
                )
            )
        db.commit()
        owner_id = owner.id
    finally:
        db.close()

    item_queries: list[str] = []

    def capture_item_queries(_conn, _cursor, statement, _parameters, _context, _executemany) -> None:
        normalized = " ".join(statement.lower().split())
        if normalized.startswith("select") and case["item_table"] in normalized:
            item_queries.append(normalized)

    event.listen(engine, "before_cursor_execute", capture_item_queries)
    try:
        with get_client() as client:
            summary_response = client.get(case["summary_path"].format(id=owner_id))
            list_response = client.get(case["list_path"])
    finally:
        event.remove(engine, "before_cursor_execute", capture_item_queries)

    assert summary_response.status_code == 200, summary_response.text
    summary = summary_response.json()
    assert summary["items"] == []
    assert summary["uploaded_count"] == 6
    assert summary["completed_count"] == 2
    assert summary["needs_review_count"] == 1
    assert summary["failed_count"] == 1
    assert summary["canceled_count"] == 1
    assert summary["running_count"] == 1
    assert summary["preprocessing_count"] == 1
    assert summary["progress"] == 4 / 6

    assert list_response.status_code == 200, list_response.text
    listed = next(item for item in list_response.json() if item["id"] == owner_id)
    assert listed["items"] == []
    assert listed["failed_count"] == 1

    assert item_queries
    assert all("count" in query and "group by" in query for query in item_queries)
    assert all("filename" not in query and "client_file_id" not in query for query in item_queries)


@pytest.mark.parametrize(
    "status",
    ["running", "paused", "completed", "completed_with_errors", "needs_review", "failed", "canceled"],
)
def test_workflow_start_rejects_non_startable_statuses(status: str) -> None:
    from app.database import SessionLocal, init_db
    from app.models import WorkflowDefinition, WorkflowRunItem

    init_db()
    db = SessionLocal()
    try:
        workflow = WorkflowDefinition(
            name=f"start_reject_{status}_workflow",
            definition_json=json.dumps({"nodes": [], "edges": []}),
        )
        db.add(workflow)
        db.flush()
        run = WorkflowRun(
            workflow_id=workflow.id,
            workflow_name=workflow.name,
            status=status,
            total_count=1,
            completed_at=datetime.utcnow() if status in {"completed", "completed_with_errors", "needs_review", "failed", "canceled"} else None,
        )
        db.add(run)
        db.flush()
        db.add(
            WorkflowRunItem(
                run_id=run.id,
                document_id=f"doc_start_reject_{status}",
                filename=f"start_reject_{status}.png",
                upload_index=0,
                status="completed" if status == "completed_with_errors" else status,
                completed_at=datetime.utcnow() if status in {"completed", "completed_with_errors", "needs_review", "failed", "canceled"} else None,
            )
        )
        db.commit()
        run_id = run.id
    finally:
        db.close()

    with get_client() as client:
        response = client.post(f"/api/workflow-runs/{run_id}/start")
        assert response.status_code == 409, response.text
        refreshed = client.get(f"/api/workflow-runs/{run_id}").json()

    assert refreshed["status"] == status


@pytest.mark.parametrize("status", ["completed", "completed_with_errors", "needs_review", "failed", "canceled"])
def test_workflow_discard_rejects_terminal_statuses(status: str) -> None:
    from app.database import SessionLocal, init_db
    from app.models import WorkflowDefinition, WorkflowRunItem

    init_db()
    db = SessionLocal()
    try:
        workflow = WorkflowDefinition(
            name=f"discard_reject_{status}_workflow",
            definition_json=json.dumps({"nodes": [], "edges": []}),
        )
        db.add(workflow)
        db.flush()
        run = WorkflowRun(
            workflow_id=workflow.id,
            workflow_name=workflow.name,
            status=status,
            total_count=1,
            completed_at=datetime.utcnow(),
        )
        db.add(run)
        db.flush()
        db.add(
            WorkflowRunItem(
                run_id=run.id,
                document_id=f"doc_discard_reject_{status}",
                filename=f"discard_reject_{status}.png",
                upload_index=0,
                status=status,
                completed_at=datetime.utcnow(),
            )
        )
        db.commit()
        run_id = run.id
    finally:
        db.close()

    with get_client() as client:
        response = client.post(f"/api/workflow-runs/{run_id}/discard")
        assert response.status_code == 409, response.text
        refreshed = client.get(f"/api/workflow-runs/{run_id}").json()

    assert refreshed["status"] == status


def test_workflow_resume_and_discard_partial_upload(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_resume_schema")
        classifier = create_document_classifier(client, name="workflow_resume_classifier")
        checklist = create_required_field_checklist(client, name="workflow_resume_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "resume_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 2})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]

        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["0:first.png:1:1"]},
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert appended.status_code == 200, appended.text
        document_id = appended.json()["items"][0]["document_id"]
        assert appended.json()["uploaded_count"] == 1
        assert appended.json()["total_count"] == 2

        resumed = client.post(f"/api/workflow-runs/{run_id}/resume")
        assert resumed.status_code == 422, resumed.text
        payload = resumed.json()["detail"]
        assert payload["uploaded_count"] == 1
        assert payload["total_count"] == 2
        assert payload["missing_count"] == 1

        paused = client.post(f"/api/workflow-runs/{run_id}/pause")
        assert paused.status_code == 200, paused.text
        assert paused.json()["status"] == "paused"

        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["1:second.png:1:1"]},
            files=[("files", ("second.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert appended.status_code == 200, appended.text
        started = client.post(f"/api/workflow-runs/{run_id}/start")
        assert started.status_code == 200, started.text
        payload = started.json()
        assert payload["status"] == "running"
        assert payload["total_count"] == 2
        assert payload["queued_count"] == 2

        discarded = client.post(f"/api/workflow-runs/{run_id}/discard")
        assert discarded.status_code == 200, discarded.text
        payload = discarded.json()
        assert payload["status"] == "canceled"
        assert len(payload["items"]) == 2
        assert payload["uploaded_count"] == 2
        assert client.get(f"/api/documents/{document_id}").status_code == 200


def test_workflow_items_preserve_upload_index_order(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_upload_index_schema")
        classifier = create_document_classifier(client, name="workflow_upload_index_classifier")
        checklist = create_required_field_checklist(client, name="workflow_upload_index_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "upload_index_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 2})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]

        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={
                "client_file_ids": ["1:z_last.png:1:1", "0:a_first.png:1:1"],
                "upload_indexes": ["1", "0"],
            },
            files=[
                ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert appended.status_code == 200, appended.text
        payload = appended.json()
        assert [item["filename"] for item in payload["items"]] == ["a_first.png", "z_last.png"]
        assert [item["upload_index"] for item in payload["items"]] == [0, 1]


def test_workflow_resume_skips_existing_upload_indexes_when_client_ids_change(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_resume_index_schema")
        classifier = create_document_classifier(client, name="workflow_resume_index_classifier")
        checklist = create_required_field_checklist(client, name="workflow_resume_index_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "resume_index_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 3})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]

        first = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["original:0", "original:1"], "upload_indexes": ["0", "1"]},
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert first.status_code == 200, first.text
        assert first.json()["uploaded_count"] == 2

        resumed = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={
                "client_file_ids": ["changed:0", "changed:1", "changed:2"],
                "upload_indexes": ["0", "1", "2"],
            },
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("third.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert resumed.status_code == 200, resumed.text
        payload = resumed.json()
        assert payload["uploaded_count"] == 3
        assert [item["upload_index"] for item in payload["items"]] == [0, 1, 2]


def test_workflow_restart_seals_missing_upload_slots(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_seal_schema")
        classifier = create_document_classifier(client, name="workflow_seal_classifier")
        checklist = create_required_field_checklist(client, name="workflow_seal_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "seal_missing_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 3})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]

        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["original:0"], "upload_indexes": ["0"]},
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert appended.status_code == 200, appended.text
        paused = client.post(f"/api/workflow-runs/{run_id}/pause")
        assert paused.status_code == 200, paused.text

        restarted = client.post(f"/api/workflow-runs/{run_id}/restart")
        assert restarted.status_code == 200, restarted.text
        payload = restarted.json()
        assert payload["id"] != run_id
        assert payload["restarted_from_run_id"] == run_id
        assert payload["status"] == "running"
        assert payload["uploaded_count"] == 3
        assert payload["failed_count"] == 2
        assert payload["queued_count"] == 1
        assert sorted(item["upload_index"] for item in payload["items"]) == [0, 1, 2]
        assert sum(1 for item in payload["items"] if item["filename"].startswith("missing_upload_")) == 2


def test_workflow_pause_preserves_uploads_and_restart_requeues(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_pause_schema")
        classifier = create_document_classifier(client, name="workflow_pause_classifier")
        checklist = create_required_field_checklist(client, name="workflow_pause_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "pause_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text
        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 2})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]
        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["0:first.png:1:1", "1:second.png:1:1"], "upload_indexes": ["0", "1"]},
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert appended.status_code == 200, appended.text
        started = client.post(f"/api/workflow-runs/{run_id}/start")
        assert started.status_code == 200, started.text

        paused = client.post(f"/api/workflow-runs/{run_id}/pause")
        assert paused.status_code == 200, paused.text
        payload = paused.json()
        assert payload["status"] == "paused"
        assert payload["uploaded_count"] == 2
        assert {item["status"] for item in payload["items"]} == {"paused"}

        restarted = client.post(f"/api/workflow-runs/{run_id}/restart")
        assert restarted.status_code == 200, restarted.text
        payload = restarted.json()
        assert payload["id"] != run_id
        assert payload["restarted_from_run_id"] == run_id
        assert payload["status"] == "running"
        assert payload["queued_count"] == 2
        assert {item["status"] for item in payload["items"]} == {"queued"}


def test_workflow_restart_can_create_new_run_with_current_workflow_without_copying_documents(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_restart_target_schema")
        classifier = create_document_classifier(client, name="workflow_restart_target_classifier")
        checklist = create_required_field_checklist(client, name="workflow_restart_target_checklist")
        workflow_one = client.post(
            "/api/workflows",
            json={"name": "workflow_one", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        workflow_two = client.post(
            "/api/workflows",
            json={"name": "workflow_two", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow_one.status_code == 200, workflow_one.text
        assert workflow_two.status_code == 200, workflow_two.text

        run_response = client.post(
            f"/api/workflows/{workflow_one.json()['id']}/runs",
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert run_response.status_code == 200, run_response.text
        source_run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
        source_document_ids = [item["document_id"] for item in source_run["items"]]
        assert source_run["started_at"] is not None
        assert source_run["completed_at"] is None

        restarted = client.post(
            f"/api/workflow-runs/{source_run['id']}/restart",
            json={"workflow_id": workflow_two.json()["id"]},
        )
        assert restarted.status_code == 200, restarted.text
        new_run = restarted.json()
        assert new_run["id"] != source_run["id"]
        assert new_run["workflow_id"] == workflow_two.json()["id"]
        assert new_run["workflow_name"] == "workflow_two"
        assert new_run["restarted_from_run_id"] == source_run["id"]
        assert [item["document_id"] for item in new_run["items"]] == source_document_ids
        assert new_run["started_at"] is not None
        assert new_run["completed_at"] is None

        replaced_source = client.get(f"/api/workflow-runs/{source_run['id']}").json()
        assert replaced_source["status"] == "canceled"
        assert replaced_source["completed_at"] is not None


def test_workflow_enqueue_creates_waiting_run_without_copying_documents(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_enqueue_schema")
        classifier = create_document_classifier(client, name="workflow_enqueue_classifier")
        checklist = create_required_field_checklist(client, name="workflow_enqueue_checklist")
        workflow_one = client.post(
            "/api/workflows",
            json={"name": "queue_workflow_one", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        workflow_two = client.post(
            "/api/workflows",
            json={"name": "queue_workflow_two", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow_one.status_code == 200, workflow_one.text
        assert workflow_two.status_code == 200, workflow_two.text

        run_response = client.post(
            f"/api/workflows/{workflow_one.json()['id']}/runs",
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert run_response.status_code == 200, run_response.text
        source_run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
        source_document_ids = [item["document_id"] for item in source_run["items"]]
        from app.database import SessionLocal

        db = SessionLocal()
        try:
            document_count_before = db.query(Document).count()
        finally:
            db.close()

        queued = client.post(
            f"/api/workflow-runs/{source_run['id']}/enqueue",
            json={"workflow_id": workflow_two.json()["id"]},
        )
        assert queued.status_code == 200, queued.text
        payload = queued.json()
        assert payload["id"] != source_run["id"]
        assert payload["workflow_id"] == workflow_two.json()["id"]
        assert payload["workflow_name"] == "queue_workflow_two"
        assert payload["status"] == "waiting"
        assert payload["workflow_run_group_id"] == source_run["id"]
        assert payload["queued_from_run_id"] == source_run["id"]
        assert payload["queue_order"] == 2
        assert [item["document_id"] for item in payload["items"]] == source_document_ids
        assert {item["status"] for item in payload["items"]} == {"queued"}
        db = SessionLocal()
        try:
            assert db.query(Document).count() == document_count_before
        finally:
            db.close()

        refreshed_source = client.get(f"/api/workflow-runs/{source_run['id']}").json()
        assert refreshed_source["workflow_run_group_id"] == source_run["id"]
        assert refreshed_source["queue_order"] == 1


def test_workflow_waiting_run_cannot_enqueue_or_start_out_of_order(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_queue_order_schema")
        classifier = create_document_classifier(client, name="workflow_queue_order_classifier")
        checklist = create_required_field_checklist(client, name="workflow_queue_order_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "queue_order_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert run_response.status_code == 200, run_response.text
        source_run_id = run_response.json()["id"]

        first_waiting = client.post(f"/api/workflow-runs/{source_run_id}/enqueue")
        second_waiting = client.post(f"/api/workflow-runs/{source_run_id}/enqueue")
        assert first_waiting.status_code == 200, first_waiting.text
        assert second_waiting.status_code == 200, second_waiting.text
        first_waiting_id = first_waiting.json()["id"]
        second_waiting_id = second_waiting.json()["id"]

        reenqueued = client.post(f"/api/workflow-runs/{first_waiting_id}/enqueue")
        assert reenqueued.status_code == 409, reenqueued.text

        blocked_by_source = client.post(f"/api/workflow-runs/{first_waiting_id}/start")
        assert blocked_by_source.status_code == 409, blocked_by_source.text
        resume_waiting = client.post(f"/api/workflow-runs/{first_waiting_id}/resume")
        assert resume_waiting.status_code == 409, resume_waiting.text
        pause_waiting = client.post(f"/api/workflow-runs/{first_waiting_id}/pause")
        assert pause_waiting.status_code == 409, pause_waiting.text
        assert client.get(f"/api/workflow-runs/{first_waiting_id}").json()["status"] == "waiting"

        from app.database import SessionLocal

        db = SessionLocal()
        try:
            source_run = db.get(WorkflowRun, source_run_id)
            assert source_run is not None
            source_run.status = "completed"
            source_run.completed_at = datetime.utcnow()
            for item in source_run.items:
                item.status = "completed"
                item.error_message = None
                item.completed_at = datetime.utcnow()
                item.result_json = json.dumps({"document_id": item.document_id, "filename": item.filename, "node_results": {}})
            db.commit()
        finally:
            db.close()

        out_of_order = client.post(f"/api/workflow-runs/{second_waiting_id}/start")
        assert out_of_order.status_code == 409, out_of_order.text

        started = client.post(f"/api/workflow-runs/{first_waiting_id}/start")
        assert started.status_code == 200, started.text
        assert started.json()["status"] == "running"


def test_workflow_finalize_starts_next_waiting_run(monkeypatch) -> None:
    from app import workflows as workflow_module
    from app.database import SessionLocal

    dispatched: list[tuple[str, int]] = []
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)
    monkeypatch.setattr("app.workflows._dispatch_workflow_run_async", lambda run_id, generation: dispatched.append((run_id, generation)))

    with get_client() as client:
        schema = create_schema(client, name="workflow_queue_advance_schema")
        classifier = create_document_classifier(client, name="workflow_queue_advance_classifier")
        checklist = create_required_field_checklist(client, name="workflow_queue_advance_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "queue_advance_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert run_response.status_code == 200, run_response.text
        source_run_id = run_response.json()["id"]
        queued = client.post(f"/api/workflow-runs/{source_run_id}/enqueue")
        assert queued.status_code == 200, queued.text
        queued_run_id = queued.json()["id"]

        db = SessionLocal()
        try:
            source_run = db.get(WorkflowRun, source_run_id)
            assert source_run is not None
            generation = source_run.execution_generation
            source_run.status = "running"
            for item in source_run.items:
                item.status = "completed"
                item.error_message = None
                item.completed_at = datetime.utcnow()
                item.result_json = json.dumps({"document_id": item.document_id, "filename": item.filename, "node_results": {}})
            db.commit()
        finally:
            db.close()

        workflow_module._finalize_workflow_run(source_run_id, generation)

        payload = client.get(f"/api/workflow-runs/{queued_run_id}").json()
        assert payload["status"] == "running"
        assert payload["queued_count"] == 2
        assert dispatched == [(queued_run_id, 1)]

        db = SessionLocal()
        try:
            queued_items = db.query(WorkflowRunItem).filter(WorkflowRunItem.run_id == queued_run_id).all()
            assert {item.execution_generation for item in queued_items} == {1}
        finally:
            db.close()


@pytest.mark.parametrize("blocked_status", ["paused", "canceled"])
def test_workflow_finalize_does_not_advance_queue_for_blocking_statuses(monkeypatch, blocked_status: str) -> None:
    from app import workflows as workflow_module
    from app.database import SessionLocal

    dispatched: list[tuple[str, int]] = []
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)
    monkeypatch.setattr("app.workflows._dispatch_workflow_run_async", lambda run_id, generation: dispatched.append((run_id, generation)))

    with get_client() as client:
        schema = create_schema(client, name=f"workflow_queue_blocked_{blocked_status}_schema")
        classifier = create_document_classifier(client, name=f"workflow_queue_blocked_{blocked_status}_classifier")
        checklist = create_required_field_checklist(client, name=f"workflow_queue_blocked_{blocked_status}_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": f"queue_blocked_{blocked_status}_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert run_response.status_code == 200, run_response.text
        source_run_id = run_response.json()["id"]
        queued = client.post(f"/api/workflow-runs/{source_run_id}/enqueue")
        assert queued.status_code == 200, queued.text
        queued_run_id = queued.json()["id"]

        db = SessionLocal()
        try:
            source_run = db.get(WorkflowRun, source_run_id)
            assert source_run is not None
            generation = source_run.execution_generation
            source_run.status = blocked_status
            for item in source_run.items:
                item.status = blocked_status
                item.error_message = f"{blocked_status} by test"
                item.completed_at = datetime.utcnow()
            db.commit()
        finally:
            db.close()

        workflow_module._finalize_workflow_run(source_run_id, generation)

        payload = client.get(f"/api/workflow-runs/{queued_run_id}").json()
        assert payload["status"] == "waiting"
        assert dispatched == []


def test_workflow_cancel_waiting_run_preserves_shared_documents(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_cancel_waiting_schema")
        classifier = create_document_classifier(client, name="workflow_cancel_waiting_classifier")
        checklist = create_required_field_checklist(client, name="workflow_cancel_waiting_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "cancel_waiting_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert run_response.status_code == 200, run_response.text
        source_run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
        source_document_id = source_run["items"][0]["document_id"]
        from app.database import SessionLocal

        db = SessionLocal()
        try:
            source_row = db.get(WorkflowRun, source_run["id"])
            assert source_row is not None
            source_row.status = "completed"
            source_row.completed_at = datetime.utcnow()
            source_row.inference_started_at = None
            for item in source_row.items:
                item.status = "completed"
                item.completed_at = source_row.completed_at
            db.commit()
        finally:
            db.close()

        queued = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued.status_code == 200, queued.text
        queued_run_id = queued.json()["id"]

        canceled = client.post(f"/api/workflow-runs/{queued_run_id}/cancel-waiting")
        assert canceled.status_code == 200, canceled.text
        payload = canceled.json()
        assert payload["status"] == "canceled"
        assert payload["items"][0]["status"] == "canceled"
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200


def test_workflow_queue_entry_delete_removes_run_without_deleting_documents(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_delete_queue_schema")
        classifier = create_document_classifier(client, name="workflow_delete_queue_classifier")
        checklist = create_required_field_checklist(client, name="workflow_delete_queue_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "delete_queue_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert run_response.status_code == 200, run_response.text
        source_run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
        source_document_id = source_run["items"][0]["document_id"]
        from app.database import SessionLocal

        db = SessionLocal()
        try:
            source_row = db.get(WorkflowRun, source_run["id"])
            assert source_row is not None
            source_row.status = "completed"
            source_row.completed_at = datetime.utcnow()
            source_row.inference_started_at = None
            for item in source_row.items:
                item.status = "completed"
                item.completed_at = source_row.completed_at
            db.commit()
        finally:
            db.close()

        queued = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued.status_code == 200, queued.text
        queued_run_id = queued.json()["id"]

        removed = client.delete(f"/api/workflow-runs/{queued_run_id}/queue-entry")
        assert removed.status_code == 200, removed.text
        assert removed.json() == {"status": "deleted", "id": queued_run_id}
        assert client.get(f"/api/workflow-runs/{queued_run_id}").status_code == 404
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200

        queued_again = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued_again.status_code == 200, queued_again.text
        queued_again_id = queued_again.json()["id"]
        canceled = client.post(f"/api/workflow-runs/{queued_again_id}/cancel-waiting")
        assert canceled.status_code == 200, canceled.text
        removed_canceled = client.delete(f"/api/workflow-runs/{queued_again_id}/queue-entry")
        assert removed_canceled.status_code == 200, removed_canceled.text
        assert client.get(f"/api/workflow-runs/{queued_again_id}").status_code == 404
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200

        queued_active = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued_active.status_code == 200, queued_active.text
        queued_active_id = queued_active.json()["id"]
        started = client.post(f"/api/workflow-runs/{queued_active_id}/start")
        assert started.status_code == 200, started.text
        removed_active = client.delete(f"/api/workflow-runs/{queued_active_id}/queue-entry")
        assert removed_active.status_code == 200, removed_active.text
        assert client.get(f"/api/workflow-runs/{queued_active_id}").status_code == 404
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200

        queued_paused = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued_paused.status_code == 200, queued_paused.text
        queued_paused_id = queued_paused.json()["id"]
        started_paused = client.post(f"/api/workflow-runs/{queued_paused_id}/start")
        assert started_paused.status_code == 200, started_paused.text
        paused = client.post(f"/api/workflow-runs/{queued_paused_id}/pause")
        assert paused.status_code == 200, paused.text
        removed_paused = client.delete(f"/api/workflow-runs/{queued_paused_id}/queue-entry")
        assert removed_paused.status_code == 200, removed_paused.text
        assert client.get(f"/api/workflow-runs/{queued_paused_id}").status_code == 404
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200


def test_workflow_restart_retries_failed_items_without_reupload(monkeypatch) -> None:
    def fake_failed_execute(db, item, graph):
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "failed",
            "error_message": "VLM_PROVIDER_REQUEST_FAILED: invalid api key",
            "branch_path": None,
            "path_node_ids": ["classifier"],
            "completed_node_ids": [],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
        }

    def fake_success_execute(db, item, graph):
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "completed",
            "error_message": None,
            "branch_path": "class:contract",
            "path_node_ids": ["classifier"],
            "completed_node_ids": ["classifier"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": {"status": "classified", "class_name": "contract"},
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_failed_execute)
        with get_client() as client:
            schema = create_schema(client, name="workflow_restart_schema")
            classifier = create_document_classifier(client, name="workflow_restart_classifier")
            checklist = create_required_field_checklist(client, name="workflow_restart_checklist")
            workflow = client.post(
                "/api/workflows",
                json={"name": "restart_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
            )
            assert workflow.status_code == 200, workflow.text
            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[
                    ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert run_response.status_code == 200, run_response.text
            run_id = run_response.json()["id"]
            failed_run = client.get(f"/api/workflow-runs/{run_id}").json()
            assert failed_run["status"] == "completed_with_errors"
            assert failed_run["failed_count"] == 2
            document_ids = [item["document_id"] for item in failed_run["items"]]

            monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_success_execute)
            restarted = client.post(f"/api/workflow-runs/{run_id}/restart")
            assert restarted.status_code == 200, restarted.text
            restarted_id = restarted.json()["id"]
            assert restarted_id != run_id
            payload = client.get(f"/api/workflow-runs/{restarted_id}").json()
            assert payload["status"] == "completed"
            assert payload["restarted_from_run_id"] == run_id
            assert payload["uploaded_count"] == 2
            assert [item["document_id"] for item in payload["items"]] == document_ids
            assert {item["status"] for item in payload["items"]} == {"completed"}

            restarted_completed = client.post(f"/api/workflow-runs/{restarted_id}/restart")
            assert restarted_completed.status_code == 200, restarted_completed.text
            second_restart_id = restarted_completed.json()["id"]
            assert second_restart_id != restarted_id
            payload = client.get(f"/api/workflow-runs/{second_restart_id}").json()
            assert payload["status"] == "completed"
            assert [item["document_id"] for item in payload["items"]] == document_ids
            assert {item["status"] for item in payload["items"]} == {"completed"}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_branch_without_fallback_exports_classification_only(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.document_modules.classify_document_with_vlm",
        lambda classes, allow_unknown, image_paths: {
            "status": "unknown",
            "class_name": None,
            "confidence": 0.2,
            "reason": "No class matched.",
            "evidence": [],
        },
    )
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_fallback_schema")
            classifier = create_document_classifier(client, name="workflow_fallback_classifier")
            checklist = create_required_field_checklist(client, name="workflow_fallback_checklist")
            definition = workflow_definition(schema["id"], classifier["id"], checklist["id"])
            definition["edges"] = [edge for edge in definition["edges"] if edge.get("sourceHandle") == "class:contract" or edge["source"] != "branch"]
            workflow = client.post("/api/workflows", json={"name": "fallback_missing", "definition": definition})
            assert workflow.status_code == 200, workflow.text
            assert workflow.json()["validation_warnings"]

            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", ("invoice.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            assert run["status"] == "completed"
            item = run["items"][0]
            assert item["status"] == "completed"
            assert item["error_message"] is None
            assert item["result"]["branch_path"] == "unknown"
            assert item["result"]["kie_values"] == {}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_branch_without_downstream_path_completes_classification_only() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_no_downstream_schema")
            classifier = create_document_classifier(client, name="workflow_no_downstream_classifier")
            checklist = create_required_field_checklist(client, name="workflow_no_downstream_checklist")
            definition = workflow_definition(schema["id"], classifier["id"], checklist["id"])
            definition["edges"] = [edge for edge in definition["edges"] if edge["source"] != "branch"]
            workflow = client.post("/api/workflows", json={"name": "classification_only", "definition": definition})
            assert workflow.status_code == 200, workflow.text
            assert workflow.json()["validation_warnings"]

            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", ("contract.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            item = run["items"][0]
            assert item["status"] == "completed"
            assert item["error_message"] is None
            assert item["result"]["classification"]["class_name"] == "contract"
            assert item["result"]["branch_path"] == "class:contract"
            assert item["result"]["kie_values"] == {}
            assert item["result"]["required_items"] == {}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_progress_payload_marks_current_and_completed_nodes(monkeypatch) -> None:
    observed_current_nodes: list[str | None] = []

    def fake_classifier(db, document_id, node):
        item = db.query(WorkflowRunItem).filter(WorkflowRunItem.document_id == document_id).one()
        observed_current_nodes.append(json.loads(item.result_json).get("current_node_kind"))
        return {
            "kind": "classifier",
            "status": "completed",
            "classification": {"status": "classified", "class_name": "contract"},
        }

    monkeypatch.setattr("app.workflows._execute_classifier_node", fake_classifier)
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_progress_schema")
            classifier = create_document_classifier(client, name="workflow_progress_classifier")
            checklist = create_required_field_checklist(client, name="workflow_progress_checklist")
            definition = workflow_definition(schema["id"], classifier["id"], checklist["id"])
            definition["edges"] = [edge for edge in definition["edges"] if edge["source"] != "branch"]
            workflow = client.post("/api/workflows", json={"name": "progress_payload", "definition": definition})
            assert workflow.status_code == 200, workflow.text

            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", ("progress.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            result = run["items"][0]["result"]
            assert observed_current_nodes == ["classifier"]
            assert result["current_node_id"] is None
            assert "classifier" in result["completed_node_ids"]
            assert "branch" in result["completed_node_ids"]
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_run_uses_workflow_worker_limit_for_blocking_work(monkeypatch) -> None:
    active = 0
    max_active = 0
    lock = threading.Lock()

    def fake_execute(db, item, graph):
        nonlocal active, max_active
        with lock:
            active += 1
            max_active = max(max_active, active)
        time.sleep(0.05)
        with lock:
            active -= 1
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "completed",
            "error_message": None,
            "branch_path": "class:contract",
            "path_node_ids": ["classifier"],
            "completed_node_ids": ["classifier"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": {"status": "classified", "class_name": "contract"},
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_execute)
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        os.environ["WORKFLOW_MAX_WORKERS"] = "2"
        os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = "8"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_parallel_schema")
            classifier = create_document_classifier(client, name="workflow_parallel_classifier")
            checklist = create_required_field_checklist(client, name="workflow_parallel_checklist")
            workflow = client.post(
                "/api/workflows",
                json={"name": "parallel_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
            )
            assert workflow.status_code == 200, workflow.text
            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", (f"doc_{index}.png", ONE_BY_ONE_PNG, "image/png")) for index in range(4)],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            assert run["status"] == "completed"
            assert max_active == 2
    finally:
        os.environ.pop("WORKFLOW_MAX_WORKERS", None)
        os.environ.pop("VLM_MAX_CONCURRENT_REQUESTS", None)
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_async_scheduler_does_not_backlog_paused_run(monkeypatch) -> None:
    from app import workflows as workflows_module
    from app.database import SessionLocal, init_db
    from app.models import WorkflowDefinition

    definition = {
        "nodes": [
            {"id": "input", "data": {"kind": "input", "label": "Input"}},
            {"id": "export", "data": {"kind": "export", "label": "Export"}},
        ],
        "edges": [{"id": "input-export", "source": "input", "target": "export"}],
    }

    init_db()
    db = SessionLocal()
    try:
        workflow = WorkflowDefinition(
            name=f"pause_backlog_workflow_{time.time_ns()}",
            definition_json=json.dumps(definition),
        )
        db.add(workflow)
        db.flush()
        run = WorkflowRun(
            workflow_id=workflow.id,
            workflow_name=workflow.name,
            workflow_definition_json=json.dumps(definition),
            status="queued",
            total_count=6,
            execution_generation=1,
        )
        db.add(run)
        db.flush()
        for index in range(6):
            document = Document(
                filename=f"pause_backlog_{index}.png",
                mime_type="image/png",
                size_bytes=1,
                page_count=1,
                storage_path=f"/tmp/pause_backlog_{index}.png",
                status="ready",
            )
            db.add(document)
            db.flush()
            db.add(
                WorkflowRunItem(
                    run_id=run.id,
                    document_id=document.id,
                    filename=document.filename,
                    upload_index=index,
                    status="queued",
                    execution_generation=1,
                )
            )
        db.commit()
        run_id = run.id
    finally:
        db.close()

    started: list[str] = []
    release: asyncio.Event | None = None

    async def fake_execute_graph(
        _item_id: str,
        _document_id: str,
        filename: str,
        _graph,
        _execution_generation: int,
    ) -> dict[str, Any]:
        started.append(filename)
        assert release is not None
        await release.wait()
        return {
            "document_id": _document_id,
            "filename": filename,
            "status": "completed",
            "error_message": None,
            "branch_path": None,
            "path_node_ids": ["export"],
            "completed_node_ids": ["export"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": None,
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    monkeypatch.setattr(workflows_module, "_execute_graph_for_item_async", fake_execute_graph)

    async def run_and_pause() -> None:
        nonlocal release
        release = asyncio.Event()
        task = asyncio.create_task(workflows_module.run_workflow_run_async(run_id, 1))
        for _ in range(100):
            if len(started) >= 2:
                break
            await asyncio.sleep(0.01)
        assert len(started) == 2
        await asyncio.sleep(0.03)
        assert len(started) == 2

        pause_db = SessionLocal()
        try:
            paused_run = pause_db.get(WorkflowRun, run_id)
            assert paused_run is not None
            paused_run.status = "paused"
            paused_run.execution_generation = 2
            now = datetime.utcnow()
            for item in paused_run.items:
                if item.status in {"queued", "running"}:
                    item.status = "paused"
                    item.error_message = "Paused by test"
                    item.completed_at = now
            pause_db.commit()
        finally:
            pause_db.close()

        release.set()
        await asyncio.wait_for(task, timeout=3)
        assert len(started) == 2

    try:
        monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
        get_settings.cache_clear()
        asyncio.run(run_and_pause())
    finally:
        os.environ.pop("WORKFLOW_MAX_WORKERS", None)
        get_settings.cache_clear()

    db = SessionLocal()
    try:
        paused_run = db.get(WorkflowRun, run_id)
        assert paused_run is not None
        assert paused_run.status == "paused"
        assert all(item.status == "paused" for item in paused_run.items)
    finally:
        db.close()


def test_workflow_parallel_item_failure_does_not_stop_other_items(monkeypatch) -> None:
    def fake_execute(db, item, graph):
        if item.filename == "doc_1.png":
            raise RuntimeError("synthetic workflow failure")
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "completed",
            "error_message": None,
            "branch_path": "class:contract",
            "path_node_ids": ["classifier"],
            "completed_node_ids": ["classifier"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": {"status": "classified", "class_name": "contract"},
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_execute)
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = "3"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_failure_schema")
            classifier = create_document_classifier(client, name="workflow_failure_classifier")
            checklist = create_required_field_checklist(client, name="workflow_failure_checklist")
            workflow = client.post(
                "/api/workflows",
                json={"name": "parallel_failure", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
            )
            assert workflow.status_code == 200, workflow.text
            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", (f"doc_{index}.png", ONE_BY_ONE_PNG, "image/png")) for index in range(3)],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            statuses = {item["filename"]: item["status"] for item in run["items"]}
            assert run["status"] == "completed_with_errors"
            assert statuses["doc_1.png"] == "failed"
            assert statuses["doc_0.png"] == "completed"
            assert statuses["doc_2.png"] == "completed"
    finally:
        os.environ.pop("VLM_MAX_CONCURRENT_REQUESTS", None)
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_retry_failed_items_requeues_only_failures(monkeypatch) -> None:
    attempts: dict[str, int] = {}

    def fake_execute(db, item, graph):
        attempts[item.filename] = attempts.get(item.filename, 0) + 1
        if item.filename == "doc_1.png" and attempts[item.filename] == 1:
            raise RuntimeError("synthetic transient failure")
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "completed",
            "error_message": None,
            "branch_path": "class:contract",
            "path_node_ids": ["classifier"],
            "completed_node_ids": ["classifier"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": {"status": "classified", "class_name": "contract"},
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_execute)
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = "2"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_retry_schema")
            classifier = create_document_classifier(client, name="workflow_retry_classifier")
            checklist = create_required_field_checklist(client, name="workflow_retry_checklist")
            workflow = client.post(
                "/api/workflows",
                json={"name": "retry_failed_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
            )
            assert workflow.status_code == 200, workflow.text
            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", (f"doc_{index}.png", ONE_BY_ONE_PNG, "image/png")) for index in range(2)],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            assert run["status"] == "completed_with_errors"
            assert run["failed_count"] == 1

            retry = client.post(f"/api/workflow-runs/{run['id']}/retry-failed")
            assert retry.status_code == 200, retry.text
            retried = client.get(f"/api/workflow-runs/{run['id']}").json()
            assert retried["status"] == "completed"
            assert retried["completed_count"] == 2
            assert retried["failed_count"] == 0
            assert attempts == {"doc_0.png": 1, "doc_1.png": 2}
    finally:
        os.environ.pop("VLM_MAX_CONCURRENT_REQUESTS", None)
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_batch_worker_exception_does_not_leave_batch_running(monkeypatch) -> None:
    from app import extraction as extraction_module

    monkeypatch.setattr("app.main.run_batch_jobs", lambda batch_id, job_ids: None)

    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client)
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            job_ids = [item["job_id"] for item in batch["items"]]
            original_run_job = extraction_module.run_extraction_job

            def flaky_run_job(job_id: str) -> None:
                if job_id == job_ids[0]:
                    raise RuntimeError("worker boom")
                original_run_job(job_id)

            monkeypatch.setattr(extraction_module, "run_extraction_job", flaky_run_job)
            extraction_module.run_batch_jobs(batch["id"], job_ids)

            loaded = client.get(f"/api/batches/{batch['id']}")
            assert loaded.status_code == 200, loaded.text
            payload = loaded.json()
            assert payload["status"] == "completed_with_errors"
            assert payload["progress"] == 1
            assert payload["failed_count"] == 1
            assert payload["completed_count"] == 1
            assert {item["status"] for item in payload["items"]} == {"completed", "failed"}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_raw_dependency_imports() -> None:
    import bleach  # noqa: F401
    import mammoth  # noqa: F401
    import openpyxl  # noqa: F401
    import pptx  # noqa: F401


def test_raw_extraction_pdf_upload() -> None:
    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            files={"file": ("sample.pdf", make_pdf_bytes(), "application/pdf")},
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed"
        assert payload["source_format"] == "pdf"
        assert payload["pdf_url"]
        assert payload["html_url"]

        pdf = client.get(payload["pdf_url"])
        assert pdf.status_code == 200
        assert pdf.headers["content-type"] == "application/pdf"

        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "Invoice No. INV-2026-001" in html_response.text

        recent = client.get("/api/raw-extractions").json()
        assert any(item["id"] == payload["id"] for item in recent)


def test_raw_extraction_pdf_upload_with_images_option() -> None:
    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            data={"include_images": "true"},
            files={"file": ("sample.pdf", make_pdf_with_image_bytes(), "application/pdf")},
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed"

        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "data:image/png;base64" in html_response.text


def test_raw_extraction_pdf_upload_includes_images_by_default() -> None:
    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            files={"file": ("sample.pdf", make_pdf_with_image_bytes(), "application/pdf")},
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed"

        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "data:image/png;base64" in html_response.text


def test_raw_extraction_pptx_upload_with_images_option(monkeypatch) -> None:
    def fake_convert(source_path, suffix, pdf_path):
        document = fitz.open()
        page = document.new_page(width=240, height=120)
        page.insert_text((24, 60), f"Preview for {source_path.name}")
        document.save(pdf_path)
        document.close()

    monkeypatch.setattr("app.raw_extractor.convert_office_to_pdf", fake_convert)

    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            data={"include_images": "true"},
            files={
                "file": (
                    "deck_with_image.pptx",
                    make_pptx_with_image_bytes(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed", payload
        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "data:image/png;base64" in html_response.text


def test_raw_extraction_office_uploads(monkeypatch) -> None:
    def fake_convert(source_path, suffix, pdf_path):
        document = fitz.open()
        page = document.new_page(width=240, height=120)
        page.insert_text((24, 60), f"Preview for {source_path.name}")
        document.save(pdf_path)
        document.close()

    monkeypatch.setattr("app.raw_extractor.convert_office_to_pdf", fake_convert)

    samples = [
        ("report.docx", make_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "Quarterly Report"),
        ("book.xlsx", make_xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "Revenue"),
        ("deck.pptx", make_pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "Roadmap"),
    ]
    with get_client() as client:
        for filename, data, mime_type, expected_text in samples:
            response = client.post(
                "/api/raw-extractions",
                files={"file": (filename, data, mime_type)},
            )
            assert response.status_code == 200, response.text
            payload = response.json()
            assert payload["status"] == "completed", payload
            assert payload["pdf_url"]
            assert payload["html_url"]
            assert client.get(payload["pdf_url"]).status_code == 200
            html_response = client.get(payload["html_url"])
            assert html_response.status_code == 200
            assert expected_text in html_response.text


def test_raw_extraction_xlsx_formula_option(monkeypatch) -> None:
    def fake_convert(source_path, suffix, pdf_path):
        document = fitz.open()
        page = document.new_page(width=240, height=120)
        page.insert_text((24, 60), f"Preview for {source_path.name}")
        document.save(pdf_path)
        document.close()

    monkeypatch.setattr("app.raw_extractor.convert_office_to_pdf", fake_convert)

    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            data={"include_formulas": "true"},
            files={
                "file": (
                    "book.xlsx",
                    make_xlsx_formula_bytes(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed", payload
        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "=SUM(B2:B2)" in html_response.text


def upload_png(client):
    response = client.post(
        "/api/documents",
        files={"file": ("invoice.png", ONE_BY_ONE_PNG, "image/png")},
    )
    assert response.status_code == 200, response.text
    return response.json()


def create_schema(client, name: str | None = None):
    global SCHEMA_COUNTER
    if name is None:
        SCHEMA_COUNTER += 1
        schema_name = "invoice_basic" if SCHEMA_COUNTER == 1 else f"invoice_basic_{SCHEMA_COUNTER}"
    else:
        schema_name = name
    response = client.post(
        "/api/schemas",
        json={
            "name": schema_name,
            "display_name": schema_name.replace("_", " ").title(),
            "fields": [
                {
                    "key_name": "invoice_number",
                    "description": "Invoice number near the top of the document. Return null if missing.",
                    "output_format": "string",
                },
                {
                    "key_name": "total_amount",
                    "description": "Final total amount including tax.",
                    "output_format": "float",
                },
            ],
        },
    )
    assert response.status_code == 200, response.text
    return response.json()


def create_document_classifier(client, name: str = "document_classifier"):
    response = client.post(
        "/api/document-classifiers",
        json={
            "name": name,
            "description": "문서를 사용자가 정의한 class 후보 중 하나로 분류합니다.",
            "allow_unknown": True,
            "classes": [
                {
                    "class_name": "contract",
                    "description": "계약 조건과 서명 또는 날인이 있는 문서",
                    "signals": ["계약", "서명", "날인"],
                },
                {
                    "class_name": "consent_form",
                    "description": "개인정보 또는 금융정보 조회 동의 여부가 있는 문서",
                    "signals": ["동의", "개인정보", "조회"],
                },
            ],
        },
    )
    assert response.status_code == 200, response.text
    return response.json()


def create_required_field_checklist(client, name: str = "required_checklist"):
    response = client.post(
        "/api/required-field-checklists",
        json={
            "name": name,
            "description": "필수 항목의 존재 여부만 확인합니다.",
            "regions": [
                {"id": "signature_region", "name": "서명 영역", "page": 1, "x": 0.55, "y": 0.55, "width": 0.35, "height": 0.25}
            ],
            "items": [
                {
                    "item_name": "성명",
                    "description": "성명이 문서에 존재하는지 확인합니다.",
                    "evidence_type": "text_or_handwriting",
                    "required": True,
                },
                {
                    "item_name": "서명",
                    "description": "서명 또는 날인이 존재하는지 확인합니다.",
                    "evidence_type": "signature_or_stamp",
                    "required": True,
                    "region_id": "signature_region",
                },
                {
                    "item_name": "체크박스",
                    "description": "필수 체크박스 표시가 존재하는지 확인합니다.",
                    "evidence_type": "checkbox",
                    "required": True,
                },
            ],
        },
    )
    assert response.status_code == 200, response.text
    return response.json()


def workflow_definition(schema_id: str, classifier_id: str, checklist_id: str):
    return {
        "nodes": [
            {"id": "input", "position": {"x": 0, "y": 120}, "data": {"kind": "input", "label": "Input"}},
            {
                "id": "classifier",
                "position": {"x": 220, "y": 120},
                "data": {"kind": "classifier", "label": "Classifier", "config": {"classifier_id": classifier_id}},
            },
            {"id": "branch", "position": {"x": 440, "y": 120}, "data": {"kind": "branch", "label": "Branch"}},
            {
                "id": "kie_contract",
                "position": {"x": 680, "y": 60},
                "data": {"kind": "kie", "label": "Contract KIE", "config": {"schema_id": schema_id}},
            },
            {
                "id": "required_contract",
                "position": {"x": 900, "y": 60},
                "data": {"kind": "required-checker", "label": "Contract Required", "config": {"checklist_id": checklist_id}},
            },
            {"id": "merge", "position": {"x": 1120, "y": 120}, "data": {"kind": "merge", "label": "Merge"}},
            {"id": "export", "position": {"x": 1340, "y": 120}, "data": {"kind": "export", "label": "Export"}},
        ],
        "edges": [
            {"id": "input-classifier", "source": "input", "target": "classifier"},
            {"id": "classifier-branch", "source": "classifier", "target": "branch"},
            {"id": "branch-contract", "source": "branch", "target": "kie_contract", "sourceHandle": "class:contract"},
            {"id": "branch-consent-form", "source": "branch", "target": "merge", "sourceHandle": "class:consent_form"},
            {"id": "branch-unknown", "source": "branch", "target": "merge", "sourceHandle": "unknown"},
            {"id": "kie-required", "source": "kie_contract", "target": "required_contract"},
            {"id": "required-merge", "source": "required_contract", "target": "merge"},
            {"id": "merge-export", "source": "merge", "target": "export"},
        ],
    }


def make_pdf_bytes() -> bytes:
    document = fitz.open()
    page = document.new_page(width=240, height=120)
    page.insert_text((24, 60), "Invoice No. INV-2026-001")
    buffer = io.BytesIO()
    document.save(buffer)
    document.close()
    return buffer.getvalue()


def make_pdf_with_image_bytes() -> bytes:
    document = fitz.open()
    page = document.new_page(width=240, height=120)
    page.insert_text((24, 28), "Document with image")
    page.insert_image(fitz.Rect(24, 40, 80, 96), stream=ONE_BY_ONE_PNG)
    buffer = io.BytesIO()
    document.save(buffer)
    document.close()
    return buffer.getvalue()


def make_docx_bytes() -> bytes:
    buffer = io.BytesIO()
    document_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr><w:r><w:t>Quarterly Report</w:t></w:r></w:p>
    <w:p><w:r><w:t>Executive summary paragraph.</w:t></w:r></w:p>
    <w:tbl>
      <w:tr><w:tc><w:p><w:r><w:t>Metric</w:t></w:r></w:p></w:tc><w:tc><w:p><w:r><w:t>Value</w:t></w:r></w:p></w:tc></w:tr>
      <w:tr><w:tc><w:p><w:r><w:t>Revenue</w:t></w:r></w:p></w:tc><w:tc><w:p><w:r><w:t>100</w:t></w:r></w:p></w:tc></w:tr>
    </w:tbl>
  </w:body>
</w:document>"""
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        archive.writestr("word/document.xml", document_xml)
    return buffer.getvalue()


def make_xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Finance"
    sheet.append(["Metric", "Value"])
    sheet.append(["Revenue", 100])
    buffer = io.BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def make_xlsx_formula_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Finance"
    sheet.append(["Metric", "Value"])
    sheet.append(["Revenue", 100])
    sheet.append(["Total", "=SUM(B2:B2)"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def make_pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Roadmap"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(6), Inches(1))
    textbox.text_frame.text = "Launch Raw Data Extractor"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def make_pptx_with_image_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Image slide"
    image_stream = io.BytesIO(ONE_BY_ONE_PNG)
    slide.shapes.add_picture(image_stream, Inches(1), Inches(1.4), width=Inches(1))
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
